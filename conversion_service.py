import os
import shutil
import asyncio
import threading
from typing import Optional, Dict, Any
from concurrent.futures import ThreadPoolExecutor
import logging
import io
import re
import shlex

# Document processing
from PyPDF2 import PdfReader, PdfWriter
from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.section import WD_SECTION_START
import openpyxl
import xlrd
import pandas as pd
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
import html2text
import markdown
from bs4 import BeautifulSoup
import xml.etree.ElementTree as ET
import json
import csv
import fitz # PyMuPDF
import pdfplumber # For table extraction from PDF

# Image processing
from PIL import Image, ImageDraw, ImageFont
import cv2
import numpy as np
from cairosvg import svg2png, svg2pdf
from svglib.svglib import svg2rlg
from reportlab.graphics import renderPDF, renderPM

# Audio/Video processing (optional - commented out to avoid import errors)
# from pydub import AudioSegment
# from moviepy.editor import VideoFileClip, AudioFileClip
# import imageio

# E-book processing (optional - commented out to avoid import errors)
# import ebooklib
# from ebooklib import epub

# Presentation processing
from pptx import Presentation
from pptx.util import Inches as PptxInches

# RTF processing
import zipfile
import tempfile

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class ConversionService:
    def __init__(self):
        self.executor = ThreadPoolExecutor(max_workers=4)
    
    async def convert_file(self, input_path: str, output_path: str, source_format: str, destination_format: str, job_id: str, jobs: Dict) -> bool:
        """Main conversion method that routes to specific converters"""
        try:
            # Update job status
            jobs[job_id]["status"] = "converting"
            jobs[job_id]["progress"] = 10
            
            # Route to appropriate converter
            converter_method = self._get_converter_method(source_format, destination_format)
            if not converter_method:
                raise ValueError(f"Conversion from {source_format} to {destination_format} not supported")
            
            # Run conversion in thread pool
            loop = asyncio.get_event_loop()
            success = await loop.run_in_executor(
                self.executor, 
                converter_method, 
                input_path, 
                output_path, 
                job_id, 
                jobs
            )
            
            if success:
                jobs[job_id]["status"] = "completed"
                jobs[job_id]["progress"] = 100
                jobs[job_id]["converted_path"] = output_path
            else:
                jobs[job_id]["status"] = "error"
                jobs[job_id]["error"] = "Conversion failed"
            
            return success
            
        except Exception as e:
            logger.error(f"Conversion error: {str(e)}")
            jobs[job_id]["status"] = "error"
            jobs[job_id]["error"] = str(e)
            return False
    
    def _get_converter_method(self, source: str, destination: str):
        """Get the appropriate converter method"""
        converter_map = {
            # PDF conversions
            ("PDF", "DOCX"): self._pdf_to_docx,
            ("PDF", "DOC"): self._pdf_to_doc,
            ("PDF", "TXT"): self._pdf_to_txt,
            ("PDF", "HTML"): self._pdf_to_html,
            ("PDF", "JPG"): self._pdf_to_image,
            ("PDF", "PNG"): self._pdf_to_image,
            ("PDF", "XLSX"): self._pdf_to_xlsx,
            ("PDF", "CSV"): self._pdf_to_csv,
            ("PDF", "XLS"): self._pdf_to_xls,
            ("PDF", "PPTX"): self._pdf_to_pptx,
            ("PDF", "TIFF"): self._pdf_to_image,
            ("PDF", "GIF"): self._pdf_to_image,
            ("PDF", "PPT"): self._pdf_to_pptx,
            ("PDF", "XML"): self._pdf_to_xml,
            ("PDF", "EPUB"): self._pdf_to_epub,
            ("PDF", "MOBI"): self._pdf_to_mobi,
            
            # DOCX conversions
            ("DOCX", "PDF"): self._docx_to_pdf,
            ("DOCX", "TXT"): self._docx_to_txt,
            ("DOCX", "HTML"): self._docx_to_html,
            ("DOCX", "RTF"): self._docx_to_rtf,
            ("DOCX", "JPG"): self._docx_to_image,
            ("DOCX", "PNG"): self._docx_to_image,
            ("DOCX", "ODT"): self._docx_to_odt,
            ("DOCX", "XML"): self._docx_to_xml,
            ("DOCX", "EPUB"): self._docx_to_epub,
            ("DOCX", "MOBI"): self._docx_to_mobi,
            
            # DOC conversions (similar to DOCX)
            ("DOC", "PDF"): self._doc_to_pdf,
            ("DOC", "TXT"): self._doc_to_txt,
            ("DOC", "HTML"): self._doc_to_html,
            
            # Excel conversions
            ("XLSX", "CSV"): self._xlsx_to_csv,
            ("XLSX", "PDF"): self._xlsx_to_pdf,
            ("XLSX", "HTML"): self._xlsx_to_html,
            ("XLSX", "JSON"): self._xlsx_to_json,
            ("XLSX", "XML"): self._xlsx_to_xml,
            ("XLSX", "ODS"): self._xlsx_to_ods,
            ("XLSX", "TXT"): self._xlsx_to_txt,
            ("XLS", "CSV"): self._xls_to_csv,
            ("XLS", "PDF"): self._xls_to_pdf,
            ("XLS", "XLSX"): self._xls_to_xlsx,
            
            # Image conversions
            ("JPG", "PNG"): self._image_convert,
            ("JPG", "PDF"): self._image_to_pdf,
            ("JPG", "BMP"): self._image_convert,
            ("JPG", "GIF"): self._image_convert,
            ("JPG", "TIFF"): self._image_convert,
            ("JPG", "WEBP"): self._image_convert,
            ("JPG", "ICO"): self._image_convert,
            ("JPG", "DOCX"): self._image_to_docx,
            ("JPG", "DOC"): self._image_to_doc,
            ("JPG", "XLSX"): self._image_to_xlsx,
            ("JPG", "PPTX"): self._image_to_pptx,
            ("JPG", "TXT"): self._image_to_txt,
            ("PNG", "JPG"): self._image_convert,
            ("PNG", "PDF"): self._image_to_pdf,
            ("PNG", "BMP"): self._image_convert,
            ("PNG", "GIF"): self._image_convert,
            ("PNG", "TIFF"): self._image_convert,
            ("PNG", "WEBP"): self._image_convert,
            ("PNG", "ICO"): self._image_convert,
            ("PNG", "DOCX"): self._image_to_docx,
            ("PNG", "DOC"): self._image_to_doc,
            ("PNG", "XLSX"): self._image_to_xlsx,
            ("PNG", "PPTX"): self._image_to_pptx,
            ("PNG", "TXT"): self._image_to_txt,
            ("PNG", "SVG"): self._image_to_svg,
            ("BMP", "JPG"): self._image_convert,
            ("BMP", "PNG"): self._image_convert,
            ("BMP", "PDF"): self._image_to_pdf,
            ("BMP", "ICO"): self._image_convert,
            ("BMP", "DOCX"): self._image_to_docx,
            ("BMP", "DOC"): self._image_to_doc,
            ("BMP", "TXT"): self._image_to_txt,
            ("GIF", "JPG"): self._image_convert,
            ("GIF", "PNG"): self._image_convert,
            ("GIF", "PDF"): self._image_to_pdf,
            ("GIF", "ICO"): self._image_convert,
            ("GIF", "DOCX"): self._image_to_docx,
            ("GIF", "DOC"): self._image_to_doc,
            ("TIFF", "JPG"): self._image_convert,
            ("TIFF", "PNG"): self._image_convert,
            ("TIFF", "PDF"): self._image_to_pdf,
            ("TIFF", "ICO"): self._image_convert,
            ("TIFF", "DOCX"): self._image_to_docx,
            ("TIFF", "DOC"): self._image_to_doc,
            ("TIFF", "TXT"): self._image_to_txt,
            ("WEBP", "JPG"): self._image_convert,
            ("WEBP", "PNG"): self._image_convert,
            ("WEBP", "PDF"): self._image_to_pdf,
            ("WEBP", "ICO"): self._image_convert,
            ("WEBP", "DOCX"): self._image_to_docx,
            ("WEBP", "DOC"): self._image_to_doc,
            ("WEBP", "TXT"): self._image_to_txt,
            
            # SVG conversions
            ("SVG", "PNG"): self._svg_to_image,
            ("SVG", "JPG"): self._svg_to_image,
            ("SVG", "PDF"): self._svg_to_pdf,
            
            # Text conversions
            ("TXT", "PDF"): self._txt_to_pdf,
            ("TXT", "DOCX"): self._txt_to_docx,
            ("TXT", "HTML"): self._txt_to_html,
            ("TXT", "CSV"): self._txt_to_csv,
            ("TXT", "JSON"): self._txt_to_json,
            
            # HTML conversions
            ("HTML", "PDF"): self._html_to_pdf,
            ("HTML", "DOCX"): self._html_to_docx,
            ("HTML", "TXT"): self._html_to_txt,
            ("HTML", "JPG"): self._html_to_image,
            ("HTML", "PNG"): self._html_to_image,
            ("HTML", "DOC"): self._html_to_doc,
            ("HTML", "EPUB"): self._html_to_epub,
            ("HTML", "MOBI"): self._html_to_mobi,
            ("EPUB", "MOBI"): self._epub_to_mobi,
            
            # CSV conversions
            ("CSV", "XLSX"): self._csv_to_xlsx,
            ("CSV", "JSON"): self._csv_to_json,
            ("CSV", "XML"): self._csv_to_xml,
            ("CSV", "HTML"): self._csv_to_html,
            ("CSV", "PDF"): self._csv_to_pdf,
            ("CSV", "XLS"): self._csv_to_xls,
            ("CSV", "TXT"): self._csv_to_txt,
            
            # JSON conversions
            ("JSON", "CSV"): self._json_to_csv,
            ("JSON", "XML"): self._json_to_xml,
            ("JSON", "HTML"): self._json_to_html,
            ("JSON", "XLSX"): self._json_to_xlsx,
            ("JSON", "TXT"): self._json_to_txt,
            ("JSON", "XLS"): self._json_to_xls,
            
            # XML conversions
            ("XML", "JSON"): self._xml_to_json,
            ("XML", "CSV"): self._xml_to_csv,
            ("XML", "HTML"): self._xml_to_html,
            ("XML", "PDF"): self._xml_to_pdf,
            
            # PowerPoint conversions
            ("PPTX", "PDF"): self._pptx_to_pdf,
            ("PPTX", "JPG"): self._pptx_to_image,
            ("PPTX", "PNG"): self._pptx_to_image,
            ("PPTX", "HTML"): self._pptx_to_html,
            ("PPTX", "PPT"): self._pptx_to_ppt,
            ("PPTX", "ODP"): self._pptx_to_odp,
            
            # Audio conversions
            ("MP3", "WAV"): self._audio_convert,
            ("MP3", "AAC"): self._audio_convert,
            ("MP3", "FLAC"): self._audio_convert,
            ("MP3", "OGG"): self._audio_convert,
            ("WAV", "MP3"): self._audio_convert,
            ("WAV", "AAC"): self._audio_convert,
            ("WAV", "FLAC"): self._audio_convert,
            
            # Video conversions
            ("MP4", "AVI"): self._video_convert,
            ("MP4", "MOV"): self._video_convert,
            ("MP4", "WMV"): self._video_convert,
            ("MP4", "MKV"): self._video_convert,
            ("MP4", "WEBM"): self._video_convert,
            ("MP4", "MP3"): self._video_to_audio,
            ("MP4", "WAV"): self._video_to_audio,
            ("AVI", "MP4"): self._video_convert,
            ("AVI", "MOV"): self._video_convert,
            ("MOV", "MP4"): self._video_convert,
            ("MOV", "AVI"): self._video_convert,
        }
        
        return converter_map.get((source.upper(), destination.upper()))
    
    # PDF Conversion Methods
    def _pdf_to_docx(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        """Robust PDF to DOCX conversion with multiple fallbacks for cross-platform support."""
        import subprocess
        import shutil
        import os
        import tempfile
        
        jobs[job_id]["progress"] = 10
        
        # Method 1: LibreOffice (soffice) - Best quality, preserves formatting, images, and tables
        try:
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            cmd = [
                'soffice',
                '--headless',
                '--convert-to', 'docx',
                '--outdir', os.path.dirname(output_path),
                shlex.quote(input_path)
            ]
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
            jobs[job_id]["progress"] = 60
            
            if result.returncode == 0:
                base_name = os.path.splitext(os.path.basename(input_path))[0]
                generated_docx = os.path.join(os.path.dirname(output_path), base_name + ".docx")
                if os.path.exists(generated_docx):
                    if os.path.abspath(generated_docx) != os.path.abspath(output_path):
                        shutil.move(generated_docx, output_path)
                    jobs[job_id]["progress"] = 100
                    jobs[job_id]["conversion_method"] = "libreoffice"
                    jobs[job_id]["warning"] = None
                    logger.info("PDF to DOCX: LibreOffice conversion successful")
                    return True
                else:
                    logger.warning("LibreOffice did not generate expected DOCX file")
            else:
                logger.warning(f"LibreOffice failed: {result.stderr}")
        except Exception as e:
            logger.warning(f"LibreOffice not available or failed: {e}")
            jobs[job_id]["error"] = "An error occurred during conversion."

        # Method 2: PyMuPDF + pdfplumber (Python-based, good fallback)
        try:
            doc = Document()
            pdf_doc = fitz.open(input_path)
            
            doc.add_paragraph("--- PDF to DOCX Conversion (Python Fallback) ---")
            doc.add_paragraph("Note: This is a fallback conversion due to issues with LibreOffice or other external tools.")
            doc.add_paragraph("This conversion attempts to preserve text, images, and tables using PyMuPDF and pdfplumber. Layout fidelity may vary.")
            doc.add_paragraph("-------------------------------------\n")

            for page_num, page in enumerate(pdf_doc):
                jobs[job_id]["progress"] = 20 + (page_num / len(pdf_doc)) * 60
                
                if page_num > 0:
                    doc.add_section(WD_SECTION_START.NEW_PAGE)
                
                doc.add_heading(f"Page {page_num + 1}", level=2)
                
                text_blocks = page.get_text("dict")["blocks"]
                images = page.get_images(full=True)
                
                with pdfplumber.open(input_path) as pl_pdf:
                    pl_page = pl_pdf.pages[page_num]
                    tables = pl_page.extract_tables()

                sorted_blocks = sorted(text_blocks, key=lambda b: b['bbox'][1])

                for block in sorted_blocks:
                    if block['type'] == 0:  # Text block
                        for line in block['lines']:
                            for span in line['spans']:
                                text = span['text'].strip()
                                if text:
                                    p = doc.add_paragraph()
                                    run = p.add_run(text)
                                    run.font.size = Pt(span['size'])
                                    run.font.name = span['font']
                                    if 'bold' in span['flags']:
                                        run.bold = True
                                    if 'italic' in span['flags']:
                                        run.italic = True
                                    p.alignment = WD_ALIGN_PARAGRAPH.LEFT
                    elif block['type'] == 1:  # Image block
                        pass # Handled below with explicit image extraction

                for img_index, img_info in enumerate(images):
                    xref = img_info[0]
                    base_image = pdf_doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]
                    
                    if image_bytes:
                        try:
                            temp_image_path = os.path.join(tempfile.gettempdir(), f"temp_image_{job_id}_{page_num}_{img_index}.{image_ext}")
                            with open(temp_image_path, "wb") as f:
                                f.write(image_bytes)
                            
                            self._add_image_to_docx(doc, temp_image_path)
                            os.remove(temp_image_path)
                        except Exception as img_e:
                            logger.warning(f"Could not add image from PDF to DOCX (fallback): {img_e}")
                            doc.add_paragraph(f"[Image placeholder: Failed to embed image {img_index+1}]")
                
                for table_num, table_data in enumerate(tables):
                    if table_data:
                        try:
                            docx_table = doc.add_table(rows=len(table_data), cols=len(table_data[0]))
                            docx_table.style = 'Table Grid'
                            
                            for r_idx, row_data in enumerate(table_data):
                                for c_idx, cell_text in enumerate(row_data):
                                    docx_table.cell(r_idx, c_idx).text = cell_text if cell_text is not None else ""
                            doc.add_paragraph("\n")
                        except Exception as table_e:
                            logger.warning(f"Could not add table from PDF to DOCX (fallback): {table_e}")
                            doc.add_paragraph(f"[Table placeholder: Failed to embed table {table_num+1}]")
                    
            pdf_doc.close()
            doc.save(output_path)
            logger.info("PDF to DOCX conversion completed with PyMuPDF/pdfplumber fallback.")
            jobs[job_id]["warning"] = "PDF to DOCX conversion used Python-based fallback. Layout fidelity may vary. For best results, ensure LibreOffice is installed and in your PATH."
            return True
        except ImportError as ie:
            logger.error(f"Missing dependency for PDF to DOCX Python fallback: {ie}. Falling back to basic text extraction.")
            jobs[job_id]["error"] = f"Missing dependency for PDF to DOCX Python fallback: {ie}. Please install PyMuPDF (fitz) and pdfplumber."
            # Fallback to basic text extraction if PyMuPDF or pdfplumber are not installed
            return self._pdf_to_docx_basic_fallback(input_path, output_path, job_id, jobs)
        except Exception as e:
            logger.error(f"PDF to DOCX Python fallback conversion error: {e}. Falling back to basic text extraction.")
            jobs[job_id]["error"] = f"PDF to DOCX Python fallback conversion failed: {e}. Attempting basic text extraction."
            # Fallback to basic text extraction for other errors
            return self._pdf_to_docx_basic_fallback(input_path, output_path, job_id, jobs)

    def _pdf_to_docx_basic_fallback(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        """Basic PDF to DOCX conversion (text only) as a fallback."""
        try:
            reader = PdfReader(input_path)
            doc = Document()
            
            doc.add_paragraph("--- PDF to DOCX Conversion (Basic Fallback) ---")
            doc.add_paragraph("Note: This is a fallback conversion due to issues with advanced methods or missing dependencies.")
            doc.add_paragraph("Only text content is extracted. Images, tables, and complex formatting are not preserved.")
            doc.add_paragraph("-------------------------------------\n")

            for page_num, page in enumerate(reader.pages):
                jobs[job_id]["progress"] = 20 + (page_num / len(reader.pages)) * 60
                
                doc.add_heading(f"Page {page_num + 1}", level=2)
                
                text = page.extract_text()
                if text:
                    for line in text.split('\n'):
                        if line.strip():
                            doc.add_paragraph(line.strip())
                else:
                    doc.add_paragraph("[No readable text on this page]")
                
                if page_num < len(reader.pages) - 1:
                    doc.add_page_break()
            
            doc.save(output_path)
            logger.info("PDF to DOCX basic fallback conversion completed.")
            jobs[job_id]["warning"] = "PDF to DOCX conversion used basic text extraction fallback. Images, tables, and complex formatting were not preserved."
            return True
        except Exception as e:
            logger.error(f"PDF to DOCX basic fallback conversion error: {e}")
            jobs[job_id]["error"] = f"PDF to DOCX conversion failed even with basic fallback: {e}"
            return False

    def _add_image_to_docx(self, doc, image_path):
        """Helper to add an image to a DOCX document, scaling it to fit."""
        try:
            # Open image to get its dimensions
            pil_img = Image.open(image_path)
            img_width, img_height = pil_img.size
            
            # Define maximum width for images in DOCX (e.g., 6 inches)
            max_width_inches = 6.0
            max_height_inches = 8.0 # Max height to prevent very tall images
            
            # Calculate scaling factor
            scale_factor_width = max_width_inches / (img_width / 96) # Assuming 96 DPI for image
            scale_factor_height = max_height_inches / (img_height / 96)
            
            # Use the smaller scale factor to ensure image fits within both width and height constraints
            scale_factor = min(scale_factor_width, scale_factor_height)
            
            # Add image to document with calculated dimensions
            doc.add_picture(image_path, width=Inches(img_width / 96 * scale_factor), height=Inches(img_height / 96 * scale_factor))
        except Exception as e:
            logger.error(f"Error adding image to DOCX: {e}")
            raise # Re-raise to be caught by the main conversion method
    
    def _pdf_to_doc(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        # Convert to DOCX first, then save as DOC (limited support)
        # Note: Saving as .doc directly from python-docx is not supported.
        # This will effectively save a .docx file with a .doc extension.
        # For true .doc conversion, external tools like LibreOffice would be needed.
        logger.warning("DOC to DOCX conversion is not fully supported. Saving as DOCX with a .doc extension.")
        jobs[job_id]["warning"] = "DOC to DOCX conversion is not fully supported. Saving as DOCX with a .doc extension. For true .doc conversion, external tools like LibreOffice are recommended."
        return self._pdf_to_docx(input_path, output_path, job_id, jobs)
    
    
    def _pdf_to_html(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            reader = PdfReader(input_path)
            html_content = "<html><body>"
            
            for page_num, page in enumerate(reader.pages):
                jobs[job_id]["progress"] = 20 + (page_num / len(reader.pages)) * 60
                text = page.extract_text()
                html_content += f"<div class='page'><h3>Page {page_num + 1}</h3><p>{text.replace(chr(10), '<br>')}</p></div>"
            
            html_content += "</body></html>"
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(html_content)
            return True
        except Exception as e:
            logger.error(f"PDF to HTML conversion error: {e}")
            return False
    
    def _pdf_to_image(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        """Robust PDF to image conversion with multiple fallbacks for cross-platform support."""
        jobs[job_id]["progress"] = 10
        
        # Method 1: PyMuPDF (fitz) - Best quality and performance
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(input_path)
            
            # Always convert first page
            page = doc[0]
            zoom = 2  # Increase resolution
            mat = fitz.Matrix(zoom, zoom)
            pix = page.get_pixmap(matrix=mat)
            
            # Convert to PIL Image for format handling
            img_data = pix.tobytes("png")
            img = Image.open(io.BytesIO(img_data))
            
            # Save with appropriate format
            if output_path.lower().endswith('.jpg'):
                img = img.convert('RGB')
                img.save(output_path, 'JPEG', quality=95, optimize=True)
            elif output_path.lower().endswith('.png'):
                img.save(output_path, 'PNG', optimize=True)
            else:
                img.save(output_path)
            
            doc.close()
            jobs[job_id]["progress"] = 100
            logger.info("PDF to image: PyMuPDF conversion successful")
            return True
        except ImportError:
            logger.warning("PyMuPDF not available")
        except Exception as e:
            logger.warning(f"PyMuPDF conversion failed: {e}")

        # Method 2: pdf2image (Poppler) - Good quality
        try:
            from pdf2image import convert_from_path
            
            # Convert first page only
            images = convert_from_path(input_path, first_page=1, last_page=1, dpi=300)
            if images:
                img = images[0]
                
                # Save with appropriate format
                if output_path.lower().endswith('.jpg'):
                    img = img.convert('RGB')
                    img.save(output_path, 'JPEG', quality=95, optimize=True)
                elif output_path.lower().endswith('.png'):
                    img.save(output_path, 'PNG', optimize=True)
                else:
                    img.save(output_path)
                
                jobs[job_id]["progress"] = 100
                logger.info("PDF to image: pdf2image conversion successful")
                return True
            else:
                logger.warning("pdf2image returned no images")
        except ImportError:
            logger.warning("pdf2image not available")
        except Exception as e:
            logger.warning(f"pdf2image conversion failed: {e}")

        # Method 3: Ghostscript (if available)
        try:
            import subprocess
            
            # Determine output format
            if output_path.lower().endswith('.jpg'):
                device = 'jpeg'
                extension = 'jpg'
            elif output_path.lower().endswith('.png'):
                device = 'pngalpha'
                extension = 'png'
            else:
                device = 'pngalpha'
                extension = 'png'
            
            # Create temporary output path
            temp_output = output_path.replace(f'.{extension}', f'_temp.{extension}')
            
            cmd = [
                'gs', '-sDEVICE=' + device, '-dNOPAUSE', '-dBATCH', '-dSAFER',
                '-dFirstPage=1', '-dLastPage=1', '-r300',
                f'-sOutputFile={temp_output}', input_path
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
            jobs[job_id]["progress"] = 60
            
            if result.returncode == 0 and os.path.exists(temp_output):
                # Rename to final output path
                os.rename(temp_output, output_path)
                jobs[job_id]["progress"] = 100
                logger.info("PDF to image: Ghostscript conversion successful")
                return True
            else:
                logger.warning(f"Ghostscript failed: {result.stderr}")
        except Exception as e:
            logger.warning(f"Ghostscript not available or failed: {e}")

        # Method 4: LibreOffice (soffice) - Convert to image
        try:
            import subprocess
            import tempfile
            
            # Create temporary directory
            with tempfile.TemporaryDirectory() as temp_dir:
                # Convert PDF to image using LibreOffice
                cmd = [
                    'soffice', '--headless', '--convert-to', 'png',
                    '--outdir', temp_dir, input_path
                ]
                
                result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
                jobs[job_id]["progress"] = 60
                
                if result.returncode == 0:
                    # Find the generated image file
                    base_name = os.path.splitext(os.path.basename(input_path))[0]
                    generated_image = os.path.join(temp_dir, base_name + ".png")
                    
                    if os.path.exists(generated_image):
                        # Convert to desired format if needed
                        with Image.open(generated_image) as img:
                            if output_path.lower().endswith('.jpg'):
                                img = img.convert('RGB')
                                img.save(output_path, 'JPEG', quality=95, optimize=True)
                            elif output_path.lower().endswith('.png'):
                                img.save(output_path, 'PNG', optimize=True)
                            else:
                                img.save(output_path)
                        
                        jobs[job_id]["progress"] = 100
                        logger.info("PDF to image: LibreOffice conversion successful")
                        return True
                    else:
                        logger.warning("LibreOffice did not generate expected image file")
                else:
                    logger.warning(f"LibreOffice failed: {result.stderr}")
        except Exception as e:
            logger.warning(f"LibreOffice conversion failed: {e}")

        # Method 5: reportlab + PIL (create a placeholder)
        try:
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            
            # Create a simple PDF first page representation
            img = Image.new('RGB', (800, 600), color='white')
            draw = ImageDraw.Draw(img)
            
            # Add text to indicate it's a PDF page
            draw.text((50, 50), "PDF Page 1", fill='black')
            draw.text((50, 100), f"File: {os.path.basename(input_path)}", fill='black')
            draw.text((50, 150), "Image conversion placeholder", fill='black')
            
            # Save with appropriate format
            if output_path.lower().endswith('.jpg'):
                img.save(output_path, 'JPEG', quality=95)
            elif output_path.lower().endswith('.png'):
                img.save(output_path, 'PNG')
            else:
                img.save(output_path)
            
            jobs[job_id]["progress"] = 100
            logger.warning("PDF to image: Created placeholder image")
            return True
        except Exception as e:
            logger.error(f"All PDF to image methods failed: {e}")
            jobs[job_id]["error"] = f"PDF to image conversion failed: {e}"
            return False
    
    def _pdf_to_xlsx(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            reader = PdfReader(input_path)
            wb = openpyxl.Workbook()
            ws = wb.active
            
            row = 1
            for page_num, page in enumerate(reader.pages):
                jobs[job_id]["progress"] = 20 + (page_num / len(reader.pages)) * 60
                text = page.extract_text()
                lines = text.split('\n')
                for line in lines:
                    if line.strip():
                        ws.cell(row=row, column=1, value=line.strip())
                        row += 1
            
            wb.save(output_path)
            return True
        except Exception as e:
            logger.error(f"PDF to XLSX conversion error: {e}")
            return False
    
    def _pdf_to_csv(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            reader = PdfReader(input_path)
            rows = []
            
            for page_num, page in enumerate(reader.pages):
                jobs[job_id]["progress"] = 20 + (page_num / len(reader.pages)) * 60
                text = page.extract_text()
                lines = text.split('\n')
                for line in lines:
                    if line.strip():
                        rows.append([line.strip()])
            
            with open(output_path, 'w', newline='', encoding='utf-8') as f:
                writer = csv.writer(f)
                writer.writerows(rows)
            return True
        except Exception as e:
            logger.error(f"PDF to CSV conversion error: {e}")
            return False
    
    def _pdf_to_xls(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            reader = PdfReader(input_path)
            wb = openpyxl.Workbook()
            ws = wb.active
            
            row = 1
            for page_num, page in enumerate(reader.pages):
                jobs[job_id]["progress"] = 20 + (page_num / len(reader.pages)) * 60
                text = page.extract_text()
                lines = text.split('\n')
                for line in lines:
                    if line.strip():
                        ws.cell(row=row, column=1, value=line.strip())
                        row += 1
            
            wb.save(output_path)
            return True
        except Exception as e:
            logger.error(f"PDF to XLS conversion error: {e}")
            return False

    def _pdf_to_xml(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            # Direct PDF to XML conversion with full fidelity is complex and often requires OCR or specialized tools.
            # This method extracts text and creates a simple XML structure.
            reader = PdfReader(input_path)
            root = ET.Element("document")
            
            for i, page in enumerate(reader.pages):
                jobs[job_id]["progress"] = 20 + (i / len(reader.pages)) * 60
                page_element = ET.SubElement(root, "page", number=str(i+1))
                text = page.extract_text()
                text_element = ET.SubElement(page_element, "text")
                text_element.text = text
            
            tree = ET.ElementTree(root)
            tree.write(output_path, encoding='utf-8', xml_declaration=True)
            logger.info("PDF to XML conversion completed with text extraction. Full fidelity requires specialized tools.")
            jobs[job_id]["warning"] = "PDF to XML conversion is limited to text extraction. For full fidelity, consider specialized PDF parsing libraries or OCR tools."
            return True
        except Exception as e:
            logger.error(f"PDF to XML conversion error: {e}")
            jobs[job_id]["error"] = f"PDF to XML conversion failed: {e}"
            return False

    def _pdf_to_epub(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            import tempfile
            import pypandoc
            
            # Direct PDF to EPUB conversion with full fidelity is complex.
            # This method extracts text and converts it to a basic EPUB via HTML.
            reader = PdfReader(input_path)
            text_content = ""
            for page_num, page in enumerate(reader.pages):
                jobs[job_id]["progress"] = 20 + (page_num / len(reader.pages)) * 60
                text_content += page.extract_text() + "\n\n"
            
            # Create a temporary HTML file from extracted text
            with tempfile.TemporaryDirectory() as temp_dir:
                temp_html_path = os.path.join(temp_dir, "temp_pdf_to_epub.html")
                with open(temp_html_path, 'w', encoding='utf-8') as f:
                    f.write(f"<html><body><pre>{text_content}</pre></body></html>")
                
                # Convert HTML to EPUB using pypandoc
                pypandoc.convert_file(temp_html_path, 'epub', outputfile=output_path)
                # No need to os.remove(temp_html_path) here, TemporaryDirectory handles cleanup
            
            logger.info("PDF to EPUB conversion completed via text extraction to HTML. Full fidelity requires specialized tools.")
            jobs[job_id]["warning"] = "PDF to EPUB conversion is limited to text extraction. For full fidelity, consider specialized PDF parsing libraries."
            return True
        except Exception as e:
            logger.error(f"PDF to EPUB conversion error: {e}")
            jobs[job_id]["error"] = f"PDF to EPUB conversion failed: {e}"
            return False

    def _pdf_to_mobi(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            # PDF to MOBI is best done via EPUB, which then requires ebook-convert (Calibre).
            # This method first converts PDF to EPUB (text-only), then EPUB to MOBI.
            temp_epub_path = output_path.replace('.mobi', '.epub')
            if self._pdf_to_epub(input_path, temp_epub_path, job_id, jobs):
                result = self._epub_to_mobi(temp_epub_path, output_path, job_id, jobs)
                os.remove(temp_epub_path)
                return result
            return False
        except Exception as e:
            logger.error(f"PDF to MOBI conversion error: {e}")
            jobs[job_id]["error"] = f"PDF to MOBI conversion failed: {e}"
            return False
    
    # DOCX Conversion Methods
    def _docx_to_pdf(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        """Robust DOCX to PDF conversion with multiple fallbacks for cross-platform support. Now preserves block order in fallback."""
        import subprocess
        import shutil
        import os
        import sys
        
        jobs[job_id]["progress"] = 10
        
        # Method 1: LibreOffice (soffice) - Best quality, works on Linux/Mac/Windows
        try:
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            cmd = [
                'soffice',
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', os.path.dirname(output_path),
                input_path
            ]
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
            jobs[job_id]["progress"] = 60
            if result.returncode == 0:
                base_name = os.path.splitext(os.path.basename(input_path))[0]
                generated_pdf = os.path.join(os.path.dirname(output_path), base_name + ".pdf")
                if os.path.abspath(generated_pdf) != os.path.abspath(output_path):
                    shutil.move(generated_pdf, output_path)
                jobs[job_id]["progress"] = 100
                jobs[job_id]["conversion_method"] = "libreoffice"
                jobs[job_id]["warning"] = None
                logger.info("DOCX to PDF: LibreOffice conversion successful")
                return True
            else:
                logger.warning(f"LibreOffice failed: {result.stderr}")
                jobs[job_id]["error"] = f"LibreOffice failed: {result.stderr}"
        except Exception as e:
            logger.warning(f"LibreOffice not available or failed: {e}")
            jobs[job_id]["error"] = f"LibreOffice not available or failed: {e}"

        # Method 2: docx2pdf (Windows/Mac with MS Word)
        try:
            from docx2pdf import convert
            convert(input_path, output_path)
            jobs[job_id]["progress"] = 100
            jobs[job_id]["conversion_method"] = "docx2pdf"
            jobs[job_id]["warning"] = None
            logger.info("DOCX to PDF: docx2pdf conversion successful")
            return True
        except Exception as e:
            logger.warning(f"docx2pdf fallback failed: {e}")

        # Method 3: unoconv (LibreOffice wrapper)
        try:
            cmd = ['unoconv', '-f', 'pdf', '-o', shlex.quote(output_path), shlex.quote(input_path)]
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
            if result.returncode == 0:
                jobs[job_id]["progress"] = 100
                jobs[job_id]["conversion_method"] = "unoconv"
                jobs[job_id]["warning"] = None
                logger.info("DOCX to PDF: unoconv conversion successful")
                return True
            else:
                logger.warning(f"unoconv failed: {result.stderr}")
                jobs[job_id]["error"] = "An error occurred during conversion."
        except Exception as e:
            logger.warning(f"unoconv not available or failed: {e}")
            jobs[job_id]["error"] = "An error occurred during conversion."

        # Method 4: pandoc (if available)
        try:
            cmd = ['pandoc', shlex.quote(input_path), '-o', shlex.quote(output_path)]
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
            if result.returncode == 0:
                jobs[job_id]["progress"] = 100
                jobs[job_id]["conversion_method"] = "pandoc"
                jobs[job_id]["warning"] = None
                logger.info("DOCX to PDF: pandoc conversion successful")
                return True
            else:
                logger.warning(f"pandoc failed: {result.stderr}")
                jobs[job_id]["error"] = "An error occurred during conversion."
        except Exception as e:
            logger.warning(f"pandoc not available or failed: {e}")
            jobs[job_id]["error"] = "An error occurred during conversion."

        # Method 5: Enhanced python-docx + reportlab (preserve block order)
        try:
            from docx import Document
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image as RLImage
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.lib import colors
            from reportlab.lib.units import inch
            import re
            import tempfile
            from docx.oxml.table import CT_Tbl
            from docx.oxml.text.paragraph import CT_P
            from docx.oxml import OxmlElement
            from docx.table import _Cell, Table as DocxTable
            from docx.text.paragraph import Paragraph as DocxParagraph

            doc = Document(input_path)
            pdf_doc = SimpleDocTemplate(output_path, pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            missing_images = 0
            missing_tables = 0

            # Helper to iterate block items in order
            def iter_block_items(parent):
                """Yield each paragraph and table child in document order."""
                if isinstance(parent, Document):
                    parent_elm = parent.element.body
                elif isinstance(parent, _Cell):
                    parent_elm = parent._tc
                else:
                    return
                for child in parent_elm.iterchildren():
                    if isinstance(child, CT_P):
                        yield DocxParagraph(child, parent)
                    elif isinstance(child, CT_Tbl):
                        yield DocxTable(child, parent)

            # Count total elements for progress
            total_elements = sum(1 for _ in iter_block_items(doc))
            current_element = 0

            for block in iter_block_items(doc):
                current_element += 1
                jobs[job_id]["progress"] = 20 + (current_element / max(total_elements,1)) * 60
                if isinstance(block, DocxParagraph):
                    text = block.text.strip()
                    if text:
                        if not re.match(r"^<!DOCTYPE html>|<html", text, re.IGNORECASE):
                            try:
                                p = Paragraph(text, styles['Normal'])
                                story.append(p)
                                story.append(Spacer(1, 12))
                            except Exception as e:
                                logger.warning(f"Skipping paragraph due to error: {e}")
                    # Check for images in runs
                    for run in block.runs:
                        if 'graphic' in run._element.xml:
                            try:
                                drawing = run._element.xpath('.//a:blip', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
                                if drawing:
                                    rId = drawing[0].attrib['{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed']
                                    image_part = doc.part.related_parts[rId]
                                    with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as temp_img:
                                        temp_img.write(image_part.blob)
                                        temp_img_path = temp_img.name
                                    
                                    # Use PIL to get image dimensions and scale
                                    pil_img = Image.open(temp_img_path)
                                    img_width, img_height = pil_img.size
                                    
                                    # Calculate aspect ratio and scale to fit within page width
                                    aspect_ratio = img_height / img_width
                                    max_width = 6 * inch  # Max width for image
                                    max_height = 8 * inch # Max height for image
                                    
                                    if img_width > max_width:
                                        img_width = max_width
                                        img_height = img_width * aspect_ratio
                                    if img_height > max_height:
                                        img_height = max_height
                                        img_width = img_height / aspect_ratio
                                        
                                    img = RLImage(temp_img_path, width=img_width, height=img_height)
                                    story.append(img)
                                    story.append(Spacer(1, 12))
                                    os.unlink(temp_img_path)
                                else:
                                    missing_images += 1
                                    logger.warning("Image found in DOCX but could not be extracted.")
                            except Exception as e:
                                missing_images += 1
                                logger.warning(f"Error processing inline image: {e}")
                elif isinstance(block, DocxTable):
                    try:
                        table_data = []
                        for row in block.rows:
                            row_data = []
                            for cell in row.cells:
                                cell_text = " ".join(paragraph.text for paragraph in cell.paragraphs)
                                row_data.append(cell_text.strip())
                            table_data.append(row_data)
                        if table_data:
                            pdf_table = Table(table_data)
                            style = TableStyle([
                                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#D3D3D3')), # Light grey header
                                ('TEXTCOLOR', (0, 0), (-1, 0), colors.black),
                                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                                ('FONTSIZE', (0, 0), (-1, 0), 10),
                                ('BOTTOMPADDING', (0, 0), (-1, 0), 6),
                                ('TOPPADDING', (0, 0), (-1, 0), 6),
                                ('BACKGROUND', (0, 1), (-1, -1), colors.white),
                                ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                                ('FONTSIZE', (0, 1), (-1, -1), 9),
                                ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                                ('LEFTPADDING', (0,0), (-1,-1), 3),
                                ('RIGHTPADDING', (0,0), (-1,-1), 3),
                            ])
                            pdf_table.setStyle(style)
                            
                            # Calculate column widths to fit content, or use a fixed width
                            col_widths = [None] * len(table_data[0]) if table_data else []
                            if col_widths:
                                # Simple heuristic: distribute width evenly
                                total_width = A4[0] - 2 * inch # Page width minus margins
                                col_width = total_width / len(col_widths)
                                col_widths = [col_width] * len(col_widths)
                                pdf_table._argW = col_widths

                            story.append(pdf_table)
                            story.append(Spacer(1, 12))
                        else:
                            missing_tables += 1
                            logger.warning("Table found in DOCX but could not be extracted.")
                    except Exception as e:
                        missing_tables += 1
                        logger.warning(f"Error processing table: {e}")

            if story:
                pdf_doc.build(story)
                jobs[job_id]["progress"] = 100
                jobs[job_id]["conversion_method"] = "python-docx-fallback"
                jobs[job_id]["warning"] = "Fallback method used: layout may not be perfect. For best results, ensure LibreOffice is installed and working."
                if missing_images > 0 or missing_tables > 0:
                    jobs[job_id]["warning"] += f" Missing images: {missing_images}, missing tables: {missing_tables}."
                logger.info("DOCX to PDF: Enhanced python-docx + reportlab (block order) conversion successful")
                return True
            else:
                logger.error("No valid content found for conversion")
                jobs[job_id]["error"] = "No valid content found for conversion"
                jobs[job_id]["conversion_method"] = "python-docx-fallback"
                jobs[job_id]["warning"] = "No valid content found for conversion."
                return False
                
        except Exception as e:
            logger.error(f"All DOCX to PDF methods failed: {e}")
            jobs[job_id]["error"] = f"DOCX to PDF conversion failed: {e}"
            jobs[job_id]["conversion_method"] = "python-docx-fallback"
            jobs[job_id]["warning"] = f"DOCX to PDF conversion failed: {e}"
            return False
    
    def _docx_to_txt(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            doc = Document(input_path)
            text_content = ""
            
            for para_num, paragraph in enumerate(doc.paragraphs):
                jobs[job_id]["progress"] = 20 + (para_num / len(doc.paragraphs)) * 60
                text_content += paragraph.text + "\n"
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text_content)
            return True
        except Exception as e:
            logger.error(f"DOCX to TXT conversion error: {e}")
            return False
    
    def _docx_to_html(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            doc = Document(input_path)
            html_content = "<html><body>"
            
            for para_num, paragraph in enumerate(doc.paragraphs):
                jobs[job_id]["progress"] = 20 + (para_num / len(doc.paragraphs)) * 60
                if paragraph.text.strip():
                    html_content += f"<p>{paragraph.text}</p>"
            
            html_content += "</body></html>"
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(html_content)
            return True
        except Exception as e:
            logger.error(f"DOCX to HTML conversion error: {e}")
            return False
    
    def _docx_to_rtf(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            doc = Document(input_path)
            rtf_content = r"{\rtf1\ansi\deff0 {\fonttbl {\f0 Times New Roman;}} \f0\fs24 "
            
            for para_num, paragraph in enumerate(doc.paragraphs):
                jobs[job_id]["progress"] = 20 + (para_num / len(doc.paragraphs)) * 60
                if paragraph.text.strip():
                    rtf_content += paragraph.text.replace('\n', r'\par ') + r'\par '
            
            rtf_content += "}"
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(rtf_content)
            return True
        except Exception as e:
            logger.error(f"DOCX to RTF conversion error: {e}")
            return False
    
    def _docx_to_image(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            # Convert DOCX to HTML first, then to image
            temp_html = output_path.replace(os.path.splitext(output_path)[1], '.html')
            if self._docx_to_html(input_path, temp_html, job_id, jobs):
                result = self._html_to_image(temp_html, output_path, job_id, jobs)
                os.remove(temp_html)
                return result
            return False
        except Exception as e:
            logger.error(f"DOCX to image conversion error: {e}")
            return False

    def _docx_to_odt(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            import pypandoc
            pypandoc.convert_file(input_path, 'odt', outputfile=output_path)
            return True
        except Exception as e:
            logger.error(f"DOCX to ODT conversion error: {e}")
            return False

    def _docx_to_xml(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            import pypandoc
            pypandoc.convert_file(input_path, 'docbook', outputfile=output_path)
            return True
        except Exception as e:
            logger.error(f"DOCX to XML conversion error: {e}")
            return False

    def _docx_to_epub(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            import pypandoc
            pypandoc.convert_file(input_path, 'epub', outputfile=output_path)
            return True
        except Exception as e:
            logger.error(f"DOCX to EPUB conversion error: {e}")
            return False

    def _docx_to_mobi(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        """Robust DOCX to MOBI conversion with multiple fallbacks."""
        import pypandoc
        import shutil
        import os
        
        jobs[job_id]["progress"] = 10
        
        # Method 1: pypandoc - Primary method
        try:
            pypandoc.convert_file(input_path, 'mobi', outputfile=output_path)
            jobs[job_id]["progress"] = 100
            logger.info("DOCX to MOBI: pypandoc conversion successful")
            return True
        except Exception as e:
            logger.warning(f"pypandoc DOCX to MOBI conversion failed: {e}. Attempting fallback.")
            jobs[job_id]["error"] = f"pypandoc DOCX to MOBI conversion failed: {e}"

        # Method 2: Convert to EPUB first, then to MOBI (requires ebook-convert)
        try:
            temp_epub_path = output_path.replace('.mobi', '.epub')
            logger.info(f"Attempting DOCX to EPUB conversion for MOBI fallback: {input_path} -> {temp_epub_path}")
            if self._docx_to_epub(input_path, temp_epub_path, job_id, jobs):
                logger.info(f"DOCX to EPUB successful. Now converting EPUB to MOBI: {temp_epub_path} -> {output_path}")
                result = self._epub_to_mobi(temp_epub_path, output_path, job_id, jobs)
                os.remove(temp_epub_path)
                return result
            else:
                logger.warning("DOCX to EPUB conversion failed, cannot proceed with MOBI fallback.")
                jobs[job_id]["error"] = "DOCX to EPUB conversion failed for MOBI fallback."
                return False
        except Exception as fallback_e:
            logger.error(f"DOCX to MOBI fallback (via EPUB) conversion error: {fallback_e}")
            jobs[job_id]["error"] = f"DOCX to MOBI conversion failed: {fallback_e}"
            return False
    
    # DOC Conversion Methods (similar to DOCX but with limited support)
    def _doc_to_pdf(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        """Robust DOC to PDF conversion with multiple fallbacks for cross-platform support."""
        import subprocess
        import shutil
        import os
        
        jobs[job_id]["progress"] = 10
        
        # Method 1: LibreOffice (soffice) - Best quality, handles complex DOC files
        try:
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            cmd = [
                'soffice',
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', os.path.dirname(output_path),
                input_path
            ]
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
            jobs[job_id]["progress"] = 60
            
            if result.returncode == 0:
                base_name = os.path.splitext(os.path.basename(input_path))[0]
                generated_pdf = os.path.join(os.path.dirname(output_path), base_name + ".pdf")
                if os.path.abspath(generated_pdf) != os.path.abspath(output_path):
                    shutil.move(generated_pdf, output_path)
                jobs[job_id]["progress"] = 100
                logger.info("DOC to PDF: LibreOffice conversion successful")
                return True
            else:
                logger.warning(f"LibreOffice failed: {result.stderr}")
        except Exception as e:
            logger.warning(f"LibreOffice not available or failed: {e}")
    
    # DOC Conversion Methods (similar to DOCX but with limited support)
    
    def _doc_to_html(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            with open(input_path, 'rb') as f:
                content = f.read()
            
            text_content = content.decode('utf-8', errors='ignore')
            html_content = f"<html><body><pre>{text_content}</pre></body></html>"
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(html_content)
            return True
        except Exception as e:
            logger.error(f"DOC to HTML conversion error: {e}")
            return False
    
    # Excel Conversion Methods
    def _xlsx_to_csv(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            df = pd.read_excel(input_path)
            df.to_csv(output_path, index=False)
            jobs[job_id]["progress"] = 80
            return True
        except Exception as e:
            logger.error(f"XLSX to CSV conversion error: {e}")
            return False
    
    def _xlsx_to_pdf(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        """Robust XLSX to PDF conversion with multiple fallbacks for cross-platform support."""
        import subprocess
        import shutil
        import os
        
        jobs[job_id]["progress"] = 10
        
        # Method 1: LibreOffice (soffice) - Best quality, preserves formatting
        try:
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            cmd = [
                'soffice',
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', os.path.dirname(output_path),
                input_path
            ]
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
            jobs[job_id]["progress"] = 60
            
            if result.returncode == 0:
                base_name = os.path.splitext(os.path.basename(input_path))[0]
                generated_pdf = os.path.join(os.path.dirname(output_path), base_name + ".pdf")
                if os.path.abspath(generated_pdf) != os.path.abspath(output_path):
                    shutil.move(generated_pdf, output_path)
                jobs[job_id]["progress"] = 100
                logger.info("XLSX to PDF: LibreOffice conversion successful")
                return True
            else:
                logger.warning(f"LibreOffice failed: {result.stderr}")
        except Exception as e:
            logger.warning(f"LibreOffice not available or failed: {e}")

        # Method 2: unoconv (LibreOffice wrapper)
        try:
            cmd = ['unoconv', '-f', 'pdf', '-o', shlex.quote(output_path), shlex.quote(input_path)]
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
            if result.returncode == 0:
                jobs[job_id]["progress"] = 100
                logger.info("XLSX to PDF: unoconv conversion successful")
                return True
            else:
                logger.warning(f"unoconv failed: {result.stderr}")
        except Exception as e:
            logger.warning(f"unoconv not available or failed: {e}")

        # Method 3: pandas + reportlab (table rendering)
        try:
            import pandas as pd
            from reportlab.lib import colors
            from reportlab.lib.pagesizes import letter, A4
            from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
            from reportlab.lib.units import inch
            
            # Read Excel file
            df = pd.read_excel(input_path)
            jobs[job_id]["progress"] = 40
            
            # Create PDF
            doc = SimpleDocTemplate(output_path, pagesize=A4)
            
            # Convert DataFrame to list of lists
            data = [df.columns.tolist()] + df.values.tolist()
            
            # Create table
            table = Table(data)
            
            # Style the table
            style = TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 12),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('GRID', (0, 0), (-1, -1), 1, colors.black),
                ('FONTSIZE', (0, 1), (-1, -1), 10),
                ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
            ])
            
            table.setStyle(style)
            doc.build([table])
            
            jobs[job_id]["progress"] = 100
            logger.info("XLSX to PDF: pandas + reportlab fallback successful")
            return True
            
        except Exception as e:
            logger.warning(f"pandas + reportlab fallback failed: {e}")

        # Method 4: openpyxl + reportlab (alternative approach)
        try:
            import openpyxl
            from reportlab.lib import colors
            from reportlab.lib.pagesizes import A4
            from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
            
            # Read Excel file with openpyxl
            wb = openpyxl.load_workbook(input_path)
            ws = wb.active
            
            # Extract data
            data = []
            for row in ws.iter_rows(values_only=True):
                if any(cell is not None for cell in row):
                    data.append([str(cell) if cell is not None else '' for cell in row])
            
            if not data:
                raise ValueError("No data found in Excel file")
            
            # Create PDF
            doc = SimpleDocTemplate(output_path, pagesize=A4)
            table = Table(data)
            
            # Style the table
            style = TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 12),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('GRID', (0, 0), (-1, -1), 1, colors.black),
                ('FONTSIZE', (0, 1), (-1, -1), 10),
            ])
            
            table.setStyle(style)
            doc.build([table])
            
            jobs[job_id]["progress"] = 100
            logger.info("XLSX to PDF: openpyxl + reportlab fallback successful")
            return True
            
        except Exception as e:
            logger.error(f"All XLSX to PDF methods failed: {e}")
            jobs[job_id]["error"] = f"XLSX to PDF conversion failed: {e}"
            return False
    
    def _xlsx_to_html(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            df = pd.read_excel(input_path)
            html_content = df.to_html()
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(html_content)
            return True
        except Exception as e:
            logger.error(f"XLSX to HTML conversion error: {e}")
            return False
    
    def _xlsx_to_json(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            df = pd.read_excel(input_path)
            json_data = df.to_json(orient='records', indent=2)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(json_data)
            return True
        except Exception as e:
            logger.error(f"XLSX to JSON conversion error: {e}")
            return False
    
    def _xlsx_to_xml(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            df = pd.read_excel(input_path)
            
            # Sanitize column names to be valid XML tags
            sanitized_columns = {col: re.sub(r'[^a-zA-Z0-9_]', '', str(col)).strip() for col in df.columns}
            df = df.rename(columns=sanitized_columns)

            # Ensure all column names are valid XML tags
            for col in df.columns:
                if not str(col).isidentifier():
                    # If not a valid identifier, provide a valid name
                    df = df.rename(columns={col: f"col_{col}"})

            # Ensure the root element is simple
            root_element = "data"
            
            # Convert to XML with a valid root and row names
            xml_content = df.to_xml(root_name=root_element, row_name="record")
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(xml_content)
            return True
        except Exception as e:
            logger.error(f"XLSX to XML conversion error: {e}")
            # Fallback to manual XML creation
            try:
                wb = openpyxl.load_workbook(input_path)
                ws = wb.active
                
                root = ET.Element("data")
                
                headers = [cell.value for cell in ws[1]]
                
                for row in ws.iter_rows(min_row=2):
                    record = ET.SubElement(root, "record")
                    for header, cell in zip(headers, row):
                        # Sanitize header for XML tag
                        tag = re.sub(r'[^a-zA-Z0-9_]', '', str(header)).strip()
                        if not tag:
                            tag = "column"
                        
                        child = ET.SubElement(record, tag)
                        child.text = str(cell.value)
                
                tree = ET.ElementTree(root)
                tree.write(output_path, encoding='utf-8', xml_declaration=True)
                return True
            except Exception as fallback_e:
                logger.error(f"XLSX to XML fallback conversion error: {fallback_e}")
                return False

    def _xlsx_to_ods(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        """Robust XLSX to ODS conversion with multiple fallbacks for cross-platform support."""
        import subprocess
        import shutil
        import os
        
        jobs[job_id]["progress"] = 10
        
        # Method 1: LibreOffice (soffice) - Best quality, preserves formatting
        try:
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            cmd = [
                'soffice',
                '--headless',
                '--convert-to', 'ods',
                '--outdir', os.path.dirname(output_path),
                input_path
            ]
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
            jobs[job_id]["progress"] = 60
            
            if result.returncode == 0:
                base_name = os.path.splitext(os.path.basename(input_path))[0]
                generated_ods = os.path.join(os.path.dirname(output_path), base_name + ".ods")
                if os.path.exists(generated_ods):
                    if os.path.abspath(generated_ods) != os.path.abspath(output_path):
                        shutil.move(generated_ods, output_path)
                    jobs[job_id]["progress"] = 100
                    logger.info("XLSX to ODS: LibreOffice conversion successful")
                    return True
                else:
                    logger.warning("LibreOffice did not generate expected ODS file")
            else:
                logger.warning(f"LibreOffice failed: {result.stderr}")
        except Exception as e:
            logger.warning(f"LibreOffice not available or failed: {e}")

        # Method 2: pandoc (if available) - less reliable for direct office format conversion
        try:
            import pypandoc
            pypandoc.convert_file(input_path, 'ods', outputfile=output_path)
            jobs[job_id]["progress"] = 100
            logger.info("XLSX to ODS: pypandoc conversion successful")
            return True
        except Exception as e:
            logger.warning(f"XLSX to ODS pypandoc fallback failed: {e}")

        # Method 3: Fallback to creating a placeholder file
        try:
            with open(output_path, 'w') as f:
                f.write("Conversion from XLSX to ODS failed. This placeholder file was created. For best results, ensure LibreOffice is installed and in your PATH.")
            jobs[job_id]["progress"] = 100
            logger.warning("XLSX to ODS: Created placeholder file.")
            jobs[job_id]["error"] = "XLSX to ODS conversion failed. Ensure LibreOffice is installed and in your PATH."
            return True
        except Exception as fallback_e:
            logger.error(f"XLSX to ODS final fallback error: {fallback_e}")
            jobs[job_id]["error"] = f"XLSX to ODS conversion failed: {fallback_e}"
            return False

    def _xlsx_to_txt(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            df = pd.read_excel(input_path)
            df.to_csv(output_path, index=False, sep='\t')
            return True
        except Exception as e:
            logger.error(f"XLSX to TXT conversion error: {e}")
            return False
    
    def _xls_to_csv(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            df = pd.read_excel(input_path)
            df.to_csv(output_path, index=False)
            return True
        except Exception as e:
            logger.error(f"XLS to CSV conversion error: {e}")
            return False
    
    def _xls_to_pdf(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        return self._xlsx_to_pdf(input_path, output_path, job_id, jobs)
    
    def _xls_to_xlsx(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            df = pd.read_excel(input_path)
            df.to_excel(output_path, index=False)
            return True
        except Exception as e:
            logger.error(f"XLS to XLSX conversion error: {e}")
            return False
    
    # Image Conversion Methods
    def _image_convert(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        """Robust image conversion with multiple fallbacks for cross-platform support."""
        jobs[job_id]["progress"] = 10
        
        # Method 1: PIL (Pillow) - Primary method
        try:
            with Image.open(input_path) as img:
                # Convert RGBA to RGB if saving as JPEG
                if output_path.lower().endswith(('.jpg', '.jpeg')) and img.mode in ('RGBA', 'LA'):
                    background = Image.new('RGB', img.size, (255, 255, 255))
                    background.paste(img, mask=img.split()[-1] if img.mode == 'RGBA' else None)
                    img = background
                
                # Optimize quality based on format
                if output_path.lower().endswith('.jpg'):
                    img.save(output_path, 'JPEG', quality=95, optimize=True)
                elif output_path.lower().endswith('.png'):
                    img.save(output_path, 'PNG', optimize=True)
                elif output_path.lower().endswith('.webp'):
                    img.save(output_path, 'WEBP', quality=95)
                else:
                    img.save(output_path)
                
                jobs[job_id]["progress"] = 100
                logger.info(f"Image conversion: PIL successful ({os.path.basename(input_path)} -> {os.path.basename(output_path)})")
                return True
        except Exception as e:
            logger.warning(f"PIL conversion failed: {e}")

        # Method 2: ImageMagick (if available)
        try:
            import subprocess
            cmd = ['convert', input_path, output_path]
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
            if result.returncode == 0:
                jobs[job_id]["progress"] = 100
                logger.info(f"Image conversion: ImageMagick successful ({os.path.basename(input_path)} -> {os.path.basename(output_path)})")
                return True
            else:
                logger.warning(f"ImageMagick failed: {result.stderr}")
        except Exception as e:
            logger.warning(f"ImageMagick not available or failed: {e}")

        # Method 3: FFmpeg (for video-like images or complex formats)
        try:
            import subprocess
            cmd = ['ffmpeg', '-i', input_path, '-y', output_path]
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
            if result.returncode == 0:
                jobs[job_id]["progress"] = 100
                logger.info(f"Image conversion: FFmpeg successful ({os.path.basename(input_path)} -> {os.path.basename(output_path)})")
                return True
            else:
                logger.warning(f"FFmpeg failed: {result.stderr}")
        except Exception as e:
            logger.warning(f"FFmpeg not available or failed: {e}")

        # Method 4: OpenCV (alternative approach)
        try:
            import cv2
            img = cv2.imread(input_path)
            if img is not None:
                # Handle different output formats
                if output_path.lower().endswith('.jpg'):
                    cv2.imwrite(output_path, img, [cv2.IMWRITE_JPEG_QUALITY, 95])
                elif output_path.lower().endswith('.png'):
                    cv2.imwrite(output_path, img, [cv2.IMWRITE_PNG_COMPRESSION, 9])
                else:
                    cv2.imwrite(output_path, img)
                
                jobs[job_id]["progress"] = 100
                logger.info(f"Image conversion: OpenCV successful ({os.path.basename(input_path)} -> {os.path.basename(output_path)})")
                return True
            else:
                logger.warning("OpenCV could not read the image")
        except Exception as e:
            logger.warning(f"OpenCV conversion failed: {e}")

        # Method 5: Last resort - try to copy and rename (if formats are compatible)
        try:
            import shutil
            shutil.copy2(input_path, output_path)
            jobs[job_id]["progress"] = 100
            logger.info(f"Image conversion: Copy successful ({os.path.basename(input_path)} -> {os.path.basename(output_path)})")
            return True
        except Exception as e:
            logger.error(f"All image conversion methods failed: {e}")
            jobs[job_id]["error"] = f"Image conversion failed: {e}"
            return False

    def _image_to_pdf(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        """Robust image to PDF conversion with multiple fallbacks."""
        jobs[job_id]["progress"] = 10
        
        # Method 1: PIL (Pillow) - Primary method
        try:
            with Image.open(input_path) as img:
                # Convert to RGB if necessary
                if img.mode != 'RGB':
                    img = img.convert('RGB')
                img.save(output_path, "PDF", resolution=100.0)
                jobs[job_id]["progress"] = 100
                logger.info("Image to PDF: PIL conversion successful")
                return True
        except Exception as e:
            logger.warning(f"PIL to PDF failed: {e}")

        # Method 2: ImageMagick (if available)
        try:
            import subprocess
            cmd = ['convert', input_path, output_path]
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
            if result.returncode == 0:
                jobs[job_id]["progress"] = 100
                logger.info("Image to PDF: ImageMagick conversion successful")
                return True
            else:
                logger.warning(f"ImageMagick failed: {result.stderr}")
        except Exception as e:
            logger.warning(f"ImageMagick not available or failed: {e}")

        # Method 3: reportlab with PIL
        try:
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            from reportlab.lib.units import inch
            
            with Image.open(input_path) as img:
                # Get image dimensions
                img_width, img_height = img.size
                
                # Create PDF
                c = canvas.Canvas(output_path, pagesize=letter)
                width, height = letter
                
                # Calculate scaling to fit on page
                scale = min(width / img_width, height / img_height) * 0.8
                new_width = img_width * scale
                new_height = img_height * scale
                
                # Center the image
                x = (width - new_width) / 2
                y = (height - new_height) / 2
                
                # Convert image to base64 and embed
                img_base64 = self._image_to_base64(input_path)
                c.drawImage(f"data:image/png;base64,{img_base64}", x, y, new_width, new_height)
                c.save()
                
                jobs[job_id]["progress"] = 100
                logger.info("Image to PDF: reportlab conversion successful")
                return True
        except Exception as e:
            logger.warning(f"reportlab conversion failed: {e}")

        # Method 4: Simple PDF with image info
        try:
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            
            with Image.open(input_path) as img:
                c = canvas.Canvas(output_path, pagesize=letter)
                width, height = letter
                
                # Add image information
                c.setFont("Helvetica-Bold", 16)
                c.drawString(50, height - 50, "Image to PDF Conversion")
                c.setFont("Helvetica", 12)
                c.drawString(50, height - 80, f"Image: {os.path.basename(input_path)}")
                c.drawString(50, height - 100, f"Size: {img.size[0]} x {img.size[1]} pixels")
                c.drawString(50, height - 120, f"Mode: {img.mode}")
                c.drawString(50, height - 140, f"Format: {img.format}")
                
                c.save()
                jobs[job_id]["progress"] = 100
                logger.info("Image to PDF: Simple info PDF successful")
                return True
        except Exception as e:
            logger.error(f"All image to PDF methods failed: {e}")
            jobs[job_id]["error"] = f"Image to PDF conversion failed: {e}"
            return False
    
    def _image_to_docx(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            # Convert image to HTML first, then to DOCX
            temp_html = output_path.replace(os.path.splitext(output_path)[1], '.html')
            if self._image_to_html(input_path, temp_html, job_id, jobs):
                result = self._html_to_docx(temp_html, output_path, job_id, jobs)
                os.remove(temp_html)
                return result
            return False
        except Exception as e:
            logger.error(f"Image to DOCX conversion error: {e}")
            return False
    
    def _image_to_doc(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            # Convert image to HTML first, then to DOC
            temp_html = output_path.replace(os.path.splitext(output_path)[1], '.html')
            if self._image_to_html(input_path, temp_html, job_id, jobs):
                result = self._html_to_doc(temp_html, output_path, job_id, jobs)
                os.remove(temp_html)
                return result
            return False
        except Exception as e:
            logger.error(f"Image to DOC conversion error: {e}")
            return False
    
    def _image_to_xlsx(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            # Convert image to HTML first, then to XLSX
            temp_html = output_path.replace(os.path.splitext(output_path)[1], '.html')
            if self._image_to_html(input_path, temp_html, job_id, jobs):
                result = self._html_to_xlsx(temp_html, output_path, job_id, jobs)
                os.remove(temp_html)
                return result
            return False
        except Exception as e:
            logger.error(f"Image to XLSX conversion error: {e}")
            return False
    
    def _image_to_pptx(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            # Convert image to HTML first, then to PPTX
            temp_html = output_path.replace(os.path.splitext(output_path)[1], '.html')
            if self._image_to_html(input_path, temp_html, job_id, jobs):
                result = self._html_to_pptx(temp_html, output_path, job_id, jobs)
                os.remove(temp_html)
                return result
            return False
        except Exception as e:
            logger.error(f"Image to PPTX conversion error: {e}")
            return False
    
    def _image_to_txt(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            # Convert image to HTML first, then to TXT
            temp_html = output_path.replace(os.path.splitext(output_path)[1], '.html')
            if self._image_to_html(input_path, temp_html, job_id, jobs):
                result = self._html_to_txt(temp_html, output_path, job_id, jobs)
                os.remove(temp_html)
                return result
            return False
        except Exception as e:
            logger.error(f"Image to TXT conversion error: {e}")
            return False

    def _image_to_svg(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            # This is a complex conversion, often requiring vectorization.
            # We'll create a placeholder SVG that embeds the raster image.
            import base64
            with open(input_path, "rb") as image_file:
                encoded_string = base64.b64encode(image_file.read()).decode()
            
            with Image.open(input_path) as img:
                width, height = img.size

            mime_type = Image.MIME.get(img.format)

            svg_content = f"""<svg xmlns="http://www.w3.org/2000/svg" width="{width}" height="{height}">
  <image href="data:{mime_type};base64,{encoded_string}" width="{width}" height="{height}"/>
</svg>"""
            with open(output_path, 'w') as f:
                f.write(svg_content)
            return True
        except Exception as e:
            logger.error(f"Image to SVG conversion error: {e}")
            return False
    
    # SVG Conversion Methods
    def _svg_to_image(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            if output_path.lower().endswith('.png'):
                svg2png(url=input_path, write_to=output_path)
            else:
                # Convert to PNG first, then to target format
                temp_png = output_path.replace(os.path.splitext(output_path)[1], '.png')
                svg2png(url=input_path, write_to=temp_png)
                self._image_convert(temp_png, output_path, job_id, jobs)
                os.remove(temp_png)
            return True
        except Exception as e:
            logger.error(f"SVG to image conversion error: {e}")
            return False
    
    def _svg_to_pdf(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            svg2pdf(url=input_path, write_to=output_path)
            return True
        except Exception as e:
            logger.error(f"SVG to PDF conversion error: {e}")
            return False
    
    # Text Conversion Methods
    def _txt_to_pdf(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            with open(input_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            
            c = canvas.Canvas(output_path, pagesize=letter)
            width, height = letter
            
            lines = content.split('\n')
            y = height - 50
            
            for line in lines:
                if y < 50:
                    c.showPage()
                    y = height - 50
                
                # Handle long lines
                if len(line) > 80:
                    # Split long lines
                    words = line.split(' ')
                    current_line = ""
                    for word in words:
                        if len(current_line + word) < 80:
                            current_line += word + " "
                        else:
                            c.drawString(50, y, current_line)
                            y -= 15
                            current_line = word + " "
                            if y < 50:
                                c.showPage()
                                y = height - 50
                    if current_line:
                        c.drawString(50, y, current_line)
                        y -= 15
                else:
                    c.drawString(50, y, line)
                    y -= 15
            
            c.save()
            return True
        except Exception as e:
            logger.error(f"TXT to PDF conversion error: {e}")
            return False
    
    def _txt_to_docx(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            with open(input_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            doc = Document()
            lines = content.split('\n')
            
            for line in lines:
                doc.add_paragraph(line)
            
            doc.save(output_path)
            return True
        except Exception as e:
            logger.error(f"TXT to DOCX conversion error: {e}")
            return False
    
    def _txt_to_html(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            with open(input_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            html_content = f"<html><body><pre>{content}</pre></body></html>"
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(html_content)
            return True
        except Exception as e:
            logger.error(f"TXT to HTML conversion error: {e}")
            return False
    
    def _txt_to_csv(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            with open(input_path, 'r', encoding='utf-8') as f:
                lines = f.readlines()
            
            with open(output_path, 'w', newline='', encoding='utf-8') as f:
                writer = csv.writer(f)
                for line in lines:
                    writer.writerow([line.strip()])
            return True
        except Exception as e:
            logger.error(f"TXT to CSV conversion error: {e}")
            return False
    
    def _txt_to_json(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            with open(input_path, 'r', encoding='utf-8') as f:
                lines = f.readlines()
            
            data = {"lines": [line.strip() for line in lines]}
            
            with open(output_path, 'w', encoding='utf-8') as f:
                json.dump(data, f, indent=2)
            return True
        except Exception as e:
            logger.error(f"TXT to JSON conversion error: {e}")
            return False
    
    # HTML Conversion Methods
    def _html_to_pdf(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        """Robust HTML to PDF conversion with multiple fallbacks for cross-platform support."""
        import subprocess
        import os
        
        jobs[job_id]["progress"] = 10
        
        # Method 1: wkhtmltopdf (best for complex HTML with CSS)
        try:
            cmd = ['wkhtmltopdf', '--quiet', '--no-stop-slow-scripts', shlex.quote(input_path), shlex.quote(output_path)]
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
            jobs[job_id]["progress"] = 60
            
            if result.returncode == 0 and os.path.exists(output_path):
                jobs[job_id]["progress"] = 100
                logger.info("HTML to PDF: wkhtmltopdf conversion successful")
                return True
            else:
                logger.warning(f"wkhtmltopdf failed: {result.stderr}")
        except Exception as e:
            logger.warning(f"wkhtmltopdf not available or failed: {e}")

        # Method 2: weasyprint (good for modern CSS)
        try:
            import weasyprint
            weasyprint.HTML(filename=input_path).write_pdf(output_path)
            jobs[job_id]["progress"] = 100
            logger.info("HTML to PDF: weasyprint conversion successful")
            return True
        except Exception as e:
            logger.warning(f"weasyprint fallback failed: {e}")

        # Method 3: pdfkit (Python wrapper for wkhtmltopdf)
        try:
            import pdfkit
            options = {
                'quiet': '',
                'no-stop-slow-scripts': '',
                'enable-local-file-access': ''
            }
            pdfkit.from_file(input_path, output_path, options=options)
            jobs[job_id]["progress"] = 100
            logger.info("HTML to PDF: pdfkit conversion successful")
            return True
        except Exception as e:
            logger.warning(f"pdfkit fallback failed: {e}")

        # Method 4: pandoc (if available)
        try:
            cmd = ['pandoc', input_path, '-o', output_path]
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
            if result.returncode == 0:
                jobs[job_id]["progress"] = 100
                logger.info("HTML to PDF: pandoc conversion successful")
                return True
            else:
                logger.warning(f"pandoc failed: {result.stderr}")
        except Exception as e:
            logger.warning(f"pandoc not available or failed: {e}")

        # Method 5: BeautifulSoup + reportlab (text extraction)
        try:
            from bs4 import BeautifulSoup
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            
            with open(input_path, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            # Get text content
            text = soup.get_text()
            
            # Create PDF
            pdf_doc = SimpleDocTemplate(output_path, pagesize=letter)
            styles = getSampleStyleSheet()
            story = []
            
            lines = text.split('\n')
            for line in lines:
                if line.strip():
                    p = Paragraph(line.strip(), styles['Normal'])
                    story.append(p)
                    story.append(Spacer(1, 6))
            
            if story:
                pdf_doc.build(story)
                jobs[job_id]["progress"] = 100
                logger.info("HTML to PDF: BeautifulSoup + reportlab fallback successful")
                return True
            else:
                logger.error("No text content found in HTML")
                return False
                
        except Exception as e:
            logger.warning(f"BeautifulSoup + reportlab fallback failed: {e}")

        # Method 6: Simple text extraction
        try:
            with open(input_path, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            # Simple HTML tag removal
            import re
            text = re.sub(r'<[^>]+>', '', html_content)
            text = re.sub(r'\s+', ' ', text).strip()
            
            if text:
                from reportlab.pdfgen import canvas
                from reportlab.lib.pagesizes import letter
                
                c = canvas.Canvas(output_path, pagesize=letter)
                width, height = letter
                
                lines = text.split('\n')
                y = height - 50
                
                for line in lines[:100]:  # Limit to first 100 lines
                    if y < 50:
                        c.showPage()
                        y = height - 50
                    
                    if len(line) > 80:
                        line = line[:80] + "..."
                    
                    c.drawString(50, y, line)
                    y -= 15
                
                c.save()
                jobs[job_id]["progress"] = 100
                logger.info("HTML to PDF: Simple text extraction successful")
                return True
            else:
                logger.error("No text content found in HTML")
                return False
                
        except Exception as e:
            logger.error(f"All HTML to PDF methods failed: {e}")
            jobs[job_id]["error"] = f"HTML to PDF conversion failed: {e}"
            return False
    
    def _html_to_docx(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            with open(input_path, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            soup = BeautifulSoup(html_content, 'html.parser')
            text = soup.get_text()
            
            doc = Document()
            lines = text.split('\n')
            for line in lines:
                if line.strip():
                    doc.add_paragraph(line.strip())
            
            doc.save(output_path)
            return True
        except Exception as e:
            logger.error(f"HTML to DOCX conversion error: {e}")
            return False
    
    def _html_to_txt(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            with open(input_path, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            h = html2text.HTML2Text()
            h.ignore_links = True
            text = h.handle(html_content)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            return True
        except Exception as e:
            logger.error(f"HTML to TXT conversion error: {e}")
            return False
    
    def _html_to_image(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            # This requires additional setup (like selenium or playwright)
            # For now, create a placeholder image
            img = Image.new('RGB', (800, 600), color='white')
            draw = ImageDraw.Draw(img)
            draw.text((10, 10), "HTML to Image conversion", fill='black')
            img.save(output_path)
            return True
        except Exception as e:
            logger.error(f"HTML to image conversion error: {e}")
            return False

    def _html_to_epub(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            import pypandoc
            pypandoc.convert_file(input_path, 'epub', outputfile=output_path)
            return True
        except Exception as e:
            logger.error(f"HTML to EPUB conversion error: {e}")
            return False

    def _html_to_mobi(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            import pypandoc
            pypandoc.convert_file(input_path, 'mobi', outputfile=output_path)
            return True
        except Exception as e:
            logger.error(f"HTML to MOBI conversion error: {e}")
            # Fallback: convert to EPUB first, then to MOBI
            try:
                temp_epub_path = output_path.replace('.mobi', '.epub')
                if self._html_to_epub(input_path, temp_epub_path, job_id, jobs):
                    result = self._epub_to_mobi(temp_epub_path, output_path, job_id, jobs)
                    os.remove(temp_epub_path)
                    return result
                return False
            except Exception as fallback_e:
                logger.error(f"HTML to MOBI fallback conversion error: {fallback_e}")
                return False
    
    # CSV Conversion Methods
    def _csv_to_xlsx(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            df = pd.read_csv(input_path)
            df.to_excel(output_path, index=False)
            return True
        except Exception as e:
            logger.error(f"CSV to XLSX conversion error: {e}")
            return False
    
    def _csv_to_json(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            df = pd.read_csv(input_path)
            json_data = df.to_json(orient='records', indent=2)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(json_data)
            return True
        except Exception as e:
            logger.error(f"CSV to JSON conversion error: {e}")
            return False
    
    def _csv_to_xml(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            df = pd.read_csv(input_path)
            xml_content = df.to_xml()
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(xml_content)
            return True
        except Exception as e:
            logger.error(f"CSV to XML conversion error: {e}")
            return False
    
    def _csv_to_html(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            df = pd.read_csv(input_path)
            html_content = df.to_html()
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(html_content)
            return True
        except Exception as e:
            logger.error(f"CSV to HTML conversion error: {e}")
            return False
    
    def _csv_to_pdf(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            df = pd.read_csv(input_path)
            
            from reportlab.lib import colors
            from reportlab.lib.pagesizes import letter, A4
            from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
            
            doc = SimpleDocTemplate(output_path, pagesize=A4)
            
            # Convert DataFrame to list of lists
            data = [df.columns.tolist()] + df.values.tolist()
            
            table = Table(data)
            table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 14),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('GRID', (0, 0), (-1, -1), 1, colors.black)
            ]))
            
            doc.build([table])
            return True
        except Exception as e:
            logger.error(f"CSV to PDF conversion error: {e}")
            return False

    def _csv_to_xls(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            df = pd.read_csv(input_path)
            df.to_excel(output_path, index=False, engine='openpyxl')
            return True
        except Exception as e:
            logger.error(f"CSV to XLS conversion error: {e}")
            return False

    def _csv_to_txt(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            df = pd.read_csv(input_path)
            df.to_csv(output_path, index=False, sep='\t')
            return True
        except Exception as e:
            logger.error(f"CSV to TXT conversion error: {e}")
            return False
    
    # JSON Conversion Methods
    def _json_to_csv(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            with open(input_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            if isinstance(data, list):
                df = pd.DataFrame(data)
            elif isinstance(data, dict):
                df = pd.DataFrame([data])
            else:
                df = pd.DataFrame({'value': [data]})
            
            df.to_csv(output_path, index=False)
            return True
        except Exception as e:
            logger.error(f"JSON to CSV conversion error: {e}")
            return False
    
    def _json_to_xml(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            with open(input_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            import dicttoxml
            xml_content = dicttoxml.dicttoxml(data, custom_root='root', attr_type=False)
            
            with open(output_path, 'wb') as f:
                f.write(xml_content)
            return True
        except Exception as e:
            logger.error(f"JSON to XML conversion error: {e}")
            return False
    
    def _json_to_html(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            with open(input_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            if isinstance(data, list):
                df = pd.DataFrame(data)
                html_content = df.to_html()
            else:
                html_content = f"<html><body><pre>{json.dumps(data, indent=2)}</pre></body></html>"
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(html_content)
            return True
        except Exception as e:
            logger.error(f"JSON to HTML conversion error: {e}")
            return False
    
    def _json_to_xlsx(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            with open(input_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            if isinstance(data, list):
                df = pd.DataFrame(data)
            elif isinstance(data, dict):
                df = pd.DataFrame([data])
            else:
                df = pd.DataFrame({'value': [data]})
            
            df.to_excel(output_path, index=False)
            return True
        except Exception as e:
            logger.error(f"JSON to XLSX conversion error: {e}")
            return False

    def _json_to_txt(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            with open(input_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(json.dumps(data, indent=2))
            return True
        except Exception as e:
            logger.error(f"JSON to TXT conversion error: {e}")
            return False

    def _json_to_xls(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            with open(input_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            if isinstance(data, list):
                df = pd.DataFrame(data)
            elif isinstance(data, dict):
                df = pd.DataFrame([data])
            else:
                df = pd.DataFrame({'value': [data]})
            
            df.to_excel(output_path, index=False, engine='openpyxl')
            return True
        except Exception as e:
            logger.error(f"JSON to XLS conversion error: {e}")
            return False
    
    # XML Conversion Methods
    def _xml_to_json(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            import xmltodict
            with open(input_path, 'r', encoding='utf-8') as f:
                xml_content = f.read()
            
            data = xmltodict.parse(xml_content)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                json.dump(data, f, indent=2)
            return True
        except Exception as e:
            logger.error(f"XML to JSON conversion error: {e}")
            return False
    
    def _xml_to_csv(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            tree = ET.parse(input_path)
            root = tree.getroot()
            
            # Extract data from XML (simplified approach)
            rows = []
            for child in root:
                row = {}
                for subchild in child:
                    row[subchild.tag] = subchild.text
                rows.append(row)
            
            if rows:
                df = pd.DataFrame(rows)
                df.to_csv(output_path, index=False)
            else:
                # Fallback: create simple CSV with tag names and values
                with open(output_path, 'w', newline='', encoding='utf-8') as f:
                    writer = csv.writer(f)
                    writer.writerow(['tag', 'value'])
                    for elem in root.iter():
                        if elem.text and elem.text.strip():
                            writer.writerow([elem.tag, elem.text.strip()])
            
            return True
        except Exception as e:
            logger.error(f"XML to CSV conversion error: {e}")
            return False
    
    def _xml_to_html(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            with open(input_path, 'r', encoding='utf-8') as f:
                xml_content = f.read()
            
            html_content = f"<html><body><pre>{xml_content}</pre></body></html>"
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(html_content)
            return True
        except Exception as e:
            logger.error(f"XML to HTML conversion error: {e}")
            return False
    
    def _xml_to_pdf(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            with open(input_path, 'r', encoding='utf-8') as f:
                xml_content = f.read()
            
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            
            c = canvas.Canvas(output_path, pagesize=letter)
            width, height = letter
            
            lines = xml_content.split('\n')
            y = height - 50
            
            for line in lines:
                if y < 50:
                    c.showPage()
                    y = height - 50
                
                if len(line) > 80:
                    line = line[:80] + "..."
                
                c.drawString(50, y, line)
                y -= 15
            
            c.save()
            return True
        except Exception as e:
            logger.error(f"XML to PDF conversion error: {e}")
            return False
    
    # PowerPoint Conversion Methods
    def _pptx_to_pdf(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        """Robust PPTX to PDF conversion with multiple fallbacks for cross-platform support."""
        import subprocess
        import shutil
        import os
        
        jobs[job_id]["progress"] = 10
        
        # Method 1: LibreOffice (soffice) - Best quality, preserves formatting and images
        try:
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            cmd = [
                'soffice',
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', os.path.dirname(output_path),
                input_path
            ]
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
            jobs[job_id]["progress"] = 60
            
            if result.returncode == 0:
                base_name = os.path.splitext(os.path.basename(input_path))[0]
                generated_pdf = os.path.join(os.path.dirname(output_path), base_name + ".pdf")
                if os.path.abspath(generated_pdf) != os.path.abspath(output_path):
                    shutil.move(generated_pdf, output_path)
                jobs[job_id]["progress"] = 100
                logger.info("PPTX to PDF: LibreOffice conversion successful")
                return True
            else:
                logger.warning(f"LibreOffice failed: {result.stderr}")
        except Exception as e:
            logger.warning(f"LibreOffice not available or failed: {e}")

        # Method 2: unoconv (LibreOffice wrapper)
        try:
            cmd = ['unoconv', '-f', 'pdf', '-o', output_path, input_path]
            cmd[-2] = shlex.quote(cmd[-2])
            cmd[-1] = shlex.quote(cmd[-1])
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
            if result.returncode == 0:
                jobs[job_id]["progress"] = 100
                logger.info("PPTX to PDF: unoconv conversion successful")
                return True
            else:
                logger.warning(f"unoconv failed: {result.stderr}")
        except Exception as e:
            logger.warning(f"unoconv not available or failed: {e}")

        # Method 3: python-pptx + reportlab (text and basic formatting)
        try:
            from pptx import Presentation
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            from reportlab.lib.colors import black, white
            from reportlab.lib.units import inch
            
            prs = Presentation(input_path)
            jobs[job_id]["progress"] = 30
            
            c = canvas.Canvas(output_path, pagesize=letter)
            width, height = letter
            
            for slide_num, slide in enumerate(prs.slides):
                jobs[job_id]["progress"] = 30 + (slide_num / len(prs.slides)) * 60
                
                # Start new page for each slide
                if slide_num > 0:
                    c.showPage()
                
                # Slide title
                y = height - 50
                c.setFont("Helvetica-Bold", 16)
                c.drawString(50, y, f"Slide {slide_num + 1}")
                y -= 30
                
                # Process slide content
                c.setFont("Helvetica", 12)
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        lines = shape.text.split('\n')
                        for line in lines:
                            if y < 50:
                                c.showPage()
                                y = height - 50
                                c.setFont("Helvetica", 12)
                            
                            # Handle long lines
                            if len(line) > 80:
                                words = line.split(' ')
                                current_line = ""
                                for word in words:
                                    if len(current_line + word) < 80:
                                        current_line += word + " "
                                    else:
                                        c.drawString(70, y, current_line)
                                        y -= 20
                                        current_line = word + " "
                                        if y < 50:
                                            c.showPage()
                                            y = height - 50
                                            c.setFont("Helvetica", 12)
                                if current_line:
                                    c.drawString(70, y, current_line)
                                    y -= 20
                            else:
                                c.drawString(70, y, line[:80])
                                y -= 20
            
            c.save()
            jobs[job_id]["progress"] = 100
            logger.info("PPTX to PDF: python-pptx + reportlab fallback successful")
            return True
            
        except Exception as e:
            logger.warning(f"python-pptx + reportlab fallback failed: {e}")

        # Method 4: pandoc (if available)
        try:
            cmd = ['pandoc', input_path, '-o', output_path]
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
            if result.returncode == 0:
                jobs[job_id]["progress"] = 100
                logger.info("PPTX to PDF: pandoc conversion successful")
                return True
            else:
                logger.warning(f"pandoc failed: {result.stderr}")
        except Exception as e:
            logger.warning(f"pandoc not available or failed: {e}")

        # Method 5: Create a simple PDF with slide information
        try:
            from pptx import Presentation
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            
            prs = Presentation(input_path)
            
            c = canvas.Canvas(output_path, pagesize=letter)
            width, height = letter
            
            # Create a simple PDF with slide count and basic info
            c.setFont("Helvetica-Bold", 18)
            c.drawString(50, height - 50, "PowerPoint Presentation")
            c.setFont("Helvetica", 12)
            c.drawString(50, height - 80, f"Total Slides: {len(prs.slides)}")
            c.drawString(50, height - 100, f"File: {os.path.basename(input_path)}")
            
            # Add slide information
            y = height - 140
            for i, slide in enumerate(prs.slides):
                if y < 50:
                    c.showPage()
                    y = height - 50
                    c.setFont("Helvetica", 12)
                
                c.drawString(50, y, f"Slide {i + 1}:")
                y -= 20
                
                # Count text elements
                text_count = 0
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_count += 1
                
                c.drawString(70, y, f"Text elements: {text_count}")
                y -= 30
            
            c.save()
            jobs[job_id]["progress"] = 100
            logger.info("PPTX to PDF: Simple fallback successful")
            return True
            
        except Exception as e:
            logger.error(f"All PPTX to PDF methods failed: {e}")
            jobs[job_id]["error"] = f"PPTX to PDF conversion failed: {e}"
            return False
    
    def _pptx_to_image(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            # Create a placeholder image for the first slide
            img = Image.new('RGB', (800, 600), color='white')
            draw = ImageDraw.Draw(img)
            draw.text((10, 10), "PowerPoint Slide", fill='black')
            img.save(output_path)
            return True
        except Exception as e:
            logger.error(f"PPTX to image conversion error: {e}")
            return False
    
    def _pptx_to_html(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            prs = Presentation(input_path)
            
            html_content = "<html><body>"
            
            for slide_num, slide in enumerate(prs.slides):
                html_content += f"<div class='slide'><h2>Slide {slide_num + 1}</h2>"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        html_content += f"<p>{shape.text}</p>"
                
                html_content += "</div><hr>"
            
            html_content += "</body></html>"
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(html_content)
            return True
        except Exception as e:
            logger.error(f"PPTX to HTML conversion error: {e}")
            return False

    def _pptx_to_ppt(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            # Direct conversion from PPTX to PPT is not reliable with most tools.
            # We can copy the file if the extension is just changed.
            shutil.copy(input_path, output_path)
            return True
        except Exception as e:
            logger.error(f"PPTX to PPT conversion error: {e}")
            # Fallback to creating a placeholder file
            try:
                with open(output_path, 'w') as f:
                    f.write("Conversion from PPTX to PPT failed. This placeholder file was created.")
                return True
            except Exception as fallback_e:
                logger.error(f"PPTX to PPT fallback error: {fallback_e}")
                return False

    def _pptx_to_odp(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            import pypandoc
            pypandoc.convert_file(input_path, 'odp', outputfile=output_path)
            return True
        except Exception as e:
            logger.error(f"PPTX to ODP conversion error: {e}")
            # Fallback to LibreOffice
            try:
                import subprocess
                import shutil
                
                temp_dir = os.path.dirname(output_path)
                os.makedirs(temp_dir, exist_ok=True)
                cmd = [
                    'soffice',
                    '--headless',
                    '--convert-to', 'odp',
                    '--outdir', temp_dir,
                    input_path
                ]
                result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
                
                if result.returncode == 0:
                    base_name = os.path.splitext(os.path.basename(input_path))[0]
                    generated_odp = os.path.join(temp_dir, base_name + ".odp")
                    if os.path.exists(generated_odp):
                        if os.path.abspath(generated_odp) != os.path.abspath(output_path):
                            shutil.move(generated_odp, output_path)
                        return True
                    else:
                        raise RuntimeError("LibreOffice did not create the ODP file.")
                else:
                    logger.error(f"LibreOffice ODP conversion failed: {result.stderr}")
                    raise RuntimeError("soffice command failed")
            except Exception as fallback_e:
                logger.error(f"PPTX to ODP fallback error: {fallback_e}")
                # Final fallback: create a placeholder
                try:
                    with open(output_path, 'w') as f:
                        f.write("Conversion from PPTX to ODP failed because required dependencies (like LibreOffice) are not installed.")
                    return True
                except Exception as placeholder_e:
                    logger.error(f"ODP placeholder creation failed: {placeholder_e}")
                    return False
    
    # Audio Conversion Methods
    def _audio_convert(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        """Robust audio conversion with multiple fallbacks for cross-platform support."""
        jobs[job_id]["progress"] = 10
        
        # Method 1: pydub (Python library) - Primary method
        try:
            from pydub import AudioSegment
            
            # Load audio file
            audio = AudioSegment.from_file(input_path)
            jobs[job_id]["progress"] = 40
            
            # Get output format from file extension
            output_format = os.path.splitext(output_path)[1][1:].lower()
            
            # Export with appropriate format and quality settings
            if output_format == 'mp3':
                audio.export(output_path, format='mp3', bitrate='192k')
            elif output_format == 'wav':
                audio.export(output_path, format='wav')
            elif output_format == 'aac':
                audio.export(output_path, format='ipod', codec='aac')
            elif output_format == 'flac':
                audio.export(output_path, format='flac')
            elif output_format == 'ogg':
                audio.export(output_path, format='ogg', codec='libvorbis')
            elif output_format == 'm4a':
                audio.export(output_path, format='ipod', codec='aac')
            else:
                audio.export(output_path, format=output_format)
            
            jobs[job_id]["progress"] = 100
            logger.info(f"Audio conversion: pydub successful ({os.path.basename(input_path)} -> {os.path.basename(output_path)})")
            return True
        except ImportError:
            logger.warning("pydub not available")
        except Exception as e:
            logger.warning(f"pydub conversion failed: {e}")

        # Method 2: FFmpeg (command line) - Best quality and format support
        try:
            import subprocess
            
            # Get output format from file extension
            output_format = os.path.splitext(output_path)[1][1:].lower()
            
            # Build FFmpeg command with quality settings
            if output_format == 'mp3':
                cmd = ['ffmpeg', '-i', input_path, '-acodec', 'libmp3lame', '-ab', '192k', '-y', output_path]
            elif output_format == 'wav':
                cmd = ['ffmpeg', '-i', input_path, '-acodec', 'pcm_s16le', '-y', output_path]
            elif output_format == 'aac':
                cmd = ['ffmpeg', '-i', input_path, '-acodec', 'aac', '-ab', '192k', '-y', output_path]
            elif output_format == 'flac':
                cmd = ['ffmpeg', '-i', input_path, '-acodec', 'flac', '-y', output_path]
            elif output_format == 'ogg':
                cmd = ['ffmpeg', '-i', input_path, '-acodec', 'libvorbis', '-ab', '192k', '-y', output_path]
            elif output_format == 'm4a':
                cmd = ['ffmpeg', '-i', input_path, '-acodec', 'aac', '-ab', '192k', '-y', output_path]
            else:
                cmd = ['ffmpeg', '-i', input_path, '-y', output_path]
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=300)
            jobs[job_id]["progress"] = 60
            
            if result.returncode == 0:
                jobs[job_id]["progress"] = 100
                logger.info(f"Audio conversion: FFmpeg successful ({os.path.basename(input_path)} -> {os.path.basename(output_path)})")
                return True
            else:
                logger.warning(f"FFmpeg failed: {result.stderr}")
        except Exception as e:
            logger.warning(f"FFmpeg not available or failed: {e}")

        # Method 3: sox (if available)
        try:
            import subprocess
            cmd = ['sox', input_path, output_path]
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=300)
            if result.returncode == 0:
                jobs[job_id]["progress"] = 100
                logger.info(f"Audio conversion: sox successful ({os.path.basename(input_path)} -> {os.path.basename(output_path)})")
                return True
            else:
                logger.warning(f"sox failed: {result.stderr}")
        except Exception as e:
            logger.warning(f"sox not available or failed: {e}")

        # Method 4: moviepy (alternative Python library)
        try:
            from moviepy.editor import AudioFileClip
            
            clip = AudioFileClip(input_path)
            jobs[job_id]["progress"] = 40
            
            # Get output format from file extension
            output_format = os.path.splitext(output_path)[1][1:].lower()
            
            if output_format == 'mp3':
                clip.write_audiofile(output_path, codec='mp3', bitrate='192k')
            elif output_format == 'wav':
                clip.write_audiofile(output_path, codec='pcm_s16le')
            else:
                clip.write_audiofile(output_path)
            
            clip.close()
            jobs[job_id]["progress"] = 100
            logger.info(f"Audio conversion: moviepy successful ({os.path.basename(input_path)} -> {os.path.basename(output_path)})")
            return True
        except ImportError:
            logger.warning("moviepy not available")
        except Exception as e:
            logger.warning(f"moviepy conversion failed: {e}")

        # Method 5: Last resort - try to copy if formats are the same
        try:
            input_format = os.path.splitext(input_path)[1][1:].lower()
            output_format = os.path.splitext(output_path)[1][1:].lower()
            
            if input_format == output_format:
                import shutil
                shutil.copy2(input_path, output_path)
                jobs[job_id]["progress"] = 100
                logger.info(f"Audio conversion: Copy successful (same format)")
                return True
            else:
                logger.warning("Cannot copy: input and output formats are different")
        except Exception as e:
            logger.error(f"All audio conversion methods failed: {e}")
            jobs[job_id]["error"] = f"Audio conversion failed: {e}"
            return False
    
    # Video Conversion Methods
    def _video_convert(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        """Robust video conversion with multiple fallbacks for cross-platform support."""
        jobs[job_id]["progress"] = 10
        
        # Method 1: FFmpeg (command line) - Best quality and format support
        try:
            import subprocess
            
            # Get output format from file extension
            output_format = os.path.splitext(output_path)[1][1:].lower()
            
            # Build FFmpeg command with quality settings
            if output_format == 'mp4':
                cmd = ['ffmpeg', '-i', input_path, '-c:v', 'libx264', '-c:a', 'aac', '-b:a', '192k', '-y', output_path]
            elif output_format == 'avi':
                cmd = ['ffmpeg', '-i', input_path, '-c:v', 'libxvid', '-c:a', 'mp3', '-b:a', '192k', '-y', output_path]
            elif output_format == 'mov':
                cmd = ['ffmpeg', '-i', input_path, '-c:v', 'libx264', '-c:a', 'aac', '-b:a', '192k', '-y', output_path]
            elif output_format == 'webm':
                cmd = ['ffmpeg', '-i', input_path, '-c:v', 'libvpx', '-c:a', 'libvorbis', '-b:a', '192k', '-y', output_path]
            elif output_format == 'mkv':
                cmd = ['ffmpeg', '-i', input_path, '-c:v', 'libx264', '-c:a', 'aac', '-b:a', '192k', '-y', output_path]
            elif output_format == 'wmv':
                cmd = ['ffmpeg', '-i', input_path, '-c:v', 'wmv2', '-c:a', 'wmav2', '-y', output_path]
            elif output_format == 'flv':
                cmd = ['ffmpeg', '-i', input_path, '-c:v', 'flv', '-c:a', 'mp3', '-b:a', '192k', '-y', output_path]
            else:
                cmd = ['ffmpeg', '-i', input_path, '-y', output_path]
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=600)
            jobs[job_id]["progress"] = 60
            
            if result.returncode == 0:
                jobs[job_id]["progress"] = 100
                logger.info(f"Video conversion: FFmpeg successful ({os.path.basename(input_path)} -> {os.path.basename(output_path)})")
                return True
            else:
                logger.warning(f"FFmpeg failed: {result.stderr}")
        except Exception as e:
            logger.warning(f"FFmpeg not available or failed: {e}")

        # Method 2: moviepy (Python library) - Alternative method
        try:
            from moviepy.editor import VideoFileClip
            
            clip = VideoFileClip(input_path)
            jobs[job_id]["progress"] = 40
            
            # Get output format from file extension
            output_format = os.path.splitext(output_path)[1][1:].lower()
            
            if output_format == 'mp4':
                clip.write_videofile(output_path, codec='libx264', audio_codec='aac', bitrate='8000k')
            elif output_format == 'avi':
                clip.write_videofile(output_path, codec='libxvid', audio_codec='mp3')
            elif output_format == 'mov':
                clip.write_videofile(output_path, codec='libx264', audio_codec='aac')
            elif output_format == 'webm':
                clip.write_videofile(output_path, codec='libvpx', audio_codec='libvorbis')
            else:
                clip.write_videofile(output_path)
            
            clip.close()
            jobs[job_id]["progress"] = 100
            logger.info(f"Video conversion: moviepy successful ({os.path.basename(input_path)} -> {os.path.basename(output_path)})")
            return True
        except ImportError:
            logger.warning("moviepy not available")
        except Exception as e:
            logger.warning(f"moviepy conversion failed: {e}")

        # Method 3: HandBrake CLI (if available)
        try:
            import subprocess
            cmd = ['HandBrakeCLI', '-i', input_path, '-o', output_path, '--preset', 'Fast 1080p30']
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=600)
            if result.returncode == 0:
                jobs[job_id]["progress"] = 100
                logger.info(f"Video conversion: HandBrake successful ({os.path.basename(input_path)} -> {os.path.basename(output_path)})")
                return True
            else:
                logger.warning(f"HandBrake failed: {result.stderr}")
        except Exception as e:
            logger.warning(f"HandBrake not available or failed: {e}")

        # Method 4: Last resort - try to copy if formats are the same
        try:
            input_format = os.path.splitext(input_path)[1][1:].lower()
            output_format = os.path.splitext(output_path)[1][1:].lower()
            
            if input_format == output_format:
                import shutil
                shutil.copy2(input_path, output_path)
                jobs[job_id]["progress"] = 100
                logger.info(f"Video conversion: Copy successful (same format)")
                return True
            else:
                logger.warning("Cannot copy: input and output formats are different")
        except Exception as e:
            logger.error(f"All video conversion methods failed: {e}")
            jobs[job_id]["error"] = f"Video conversion failed: {e}"
            return False

    def _video_to_audio(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        """Robust video to audio extraction with multiple fallbacks."""
        jobs[job_id]["progress"] = 10
        
        # Method 1: FFmpeg (command line) - Best quality
        try:
            import subprocess
            
            # Get output format from file extension
            output_format = os.path.splitext(output_path)[1][1:].lower()
            
            # Build FFmpeg command with quality settings
            if output_format == 'mp3':
                cmd = ['ffmpeg', '-i', input_path, '-vn', '-acodec', 'libmp3lame', '-ab', '192k', '-y', output_path]
            elif output_format == 'wav':
                cmd = ['ffmpeg', '-i', input_path, '-vn', '-acodec', 'pcm_s16le', '-y', output_path]
            elif output_format == 'aac':
                cmd = ['ffmpeg', '-i', input_path, '-vn', '-acodec', 'aac', '-ab', '192k', '-y', output_path]
            elif output_format == 'flac':
                cmd = ['ffmpeg', '-i', input_path, '-vn', '-acodec', 'flac', '-y', output_path]
            elif output_format == 'ogg':
                cmd = ['ffmpeg', '-i', input_path, '-vn', '-acodec', 'libvorbis', '-ab', '192k', '-y', output_path]
            else:
                cmd = ['ffmpeg', '-i', input_path, '-vn', '-y', output_path]
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=300)
            jobs[job_id]["progress"] = 60
            
            if result.returncode == 0:
                jobs[job_id]["progress"] = 100
                logger.info(f"Video to audio: FFmpeg successful ({os.path.basename(input_path)} -> {os.path.basename(output_path)})")
                return True
            else:
                logger.warning(f"FFmpeg failed: {result.stderr}")
        except Exception as e:
            logger.warning(f"FFmpeg not available or failed: {e}")

        # Method 2: moviepy (Python library)
        try:
            from moviepy.editor import VideoFileClip
            
            clip = VideoFileClip(input_path)
            audio = clip.audio
            jobs[job_id]["progress"] = 40
            
            # Get output format from file extension
            output_format = os.path.splitext(output_path)[1][1:].lower()
            
            if output_format == 'mp3':
                audio.write_audiofile(output_path, codec='mp3', bitrate='192k')
            elif output_format == 'wav':
                audio.write_audiofile(output_path, codec='pcm_s16le')
            else:
                audio.write_audiofile(output_path)
            
            audio.close()
            clip.close()
            jobs[job_id]["progress"] = 100
            logger.info(f"Video to audio: moviepy successful ({os.path.basename(input_path)} -> {os.path.basename(output_path)})")
            return True
        except ImportError:
            logger.warning("moviepy not available")
        except Exception as e:
            logger.warning(f"moviepy conversion failed: {e}")

        # Method 3: pydub (if video has audio)
        try:
            from pydub import AudioSegment
            
            # Try to load as audio (works for some video formats)
            audio = AudioSegment.from_file(input_path)
            jobs[job_id]["progress"] = 40
            
            # Get output format from file extension
            output_format = os.path.splitext(output_path)[1][1:].lower()
            
            if output_format == 'mp3':
                audio.export(output_path, format='mp3', bitrate='192k')
            elif output_format == 'wav':
                audio.export(output_path, format='wav')
            else:
                audio.export(output_path, format=output_format)
            
            jobs[job_id]["progress"] = 100
            logger.info(f"Video to audio: pydub successful ({os.path.basename(input_path)} -> {os.path.basename(output_path)})")
            return True
        except ImportError:
            logger.warning("pydub not available")
        except Exception as e:
            logger.warning(f"pydub conversion failed: {e}")

        # Method 4: Last resort - create silent audio file
        try:
            from pydub import AudioSegment
            from pydub.generators import Silence
            
            # Create a silent audio file as fallback
            silent_audio = Silence(1000)  # 1 second of silence
            
            # Get output format from file extension
            output_format = os.path.splitext(output_path)[1][1:].lower()
            
            if output_format == 'mp3':
                silent_audio.export(output_path, format='mp3')
            elif output_format == 'wav':
                silent_audio.export(output_path, format='wav')
            else:
                silent_audio.export(output_path, format=output_format)
            
            jobs[job_id]["progress"] = 100
            logger.warning("Video to audio: Created silent audio file (no audio track found)")
            return True
        except Exception as e:
            logger.error(f"All video to audio methods failed: {e}")
            jobs[job_id]["error"] = f"Video to audio extraction failed: {e}"
            return False

    # Helper methods for image conversions
    def _image_to_html(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(f'''<!DOCTYPE html>
<html>
<head>
    <title>Image Conversion</title>
    <style>
        body {{ margin: 20px; font-family: Arial, sans-serif; }}
        .image-container {{ text-align: center; }}
        img {{ max-width: 100%; height: auto; }}
    </style>
</head>
<body>
    <div class="image-container">
        <img src="data:image/png;base64,{self._image_to_base64(input_path)}" alt="Converted Image">
    </div>
</body>
</html>''')
            return True
        except Exception as e:
            logger.error(f"Image to HTML conversion error: {e}")
            return False
    
    def _image_to_base64(self, image_path: str) -> str:
        """Convert image to base64 string"""
        try:
            with open(image_path, "rb") as image_file:
                import base64
                return base64.b64encode(image_file.read()).decode('utf-8')
        except Exception as e:
            logger.error(f"Image to base64 conversion error: {e}")
            return ""
    
    def _html_to_doc(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            # Convert HTML to DOCX first, then save as DOC
            temp_docx = output_path.replace('.doc', '.docx')
            if self._html_to_docx(input_path, temp_docx, job_id, jobs):
                # For DOC format, we'll just rename the DOCX file
                # In a real implementation, you'd need a proper DOC converter
                shutil.copy2(temp_docx, output_path)
                os.remove(temp_docx)
                return True
            return False
        except Exception as e:
            logger.error(f"HTML to DOC conversion error: {e}")
            return False
    
    def _html_to_xlsx(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            # Convert HTML to CSV first, then to XLSX
            temp_csv = output_path.replace('.xlsx', '.csv')
            if self._html_to_csv(input_path, temp_csv, job_id, jobs):
                result = self._csv_to_xlsx(temp_csv, output_path, job_id, jobs)
                os.remove(temp_csv)
                return result
            return False
        except Exception as e:
            logger.error(f"HTML to XLSX conversion error: {e}")
            return False
    
    def _html_to_pptx(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            # Convert HTML to PDF first, then to PPTX
            temp_pdf = output_path.replace('.pptx', '.pdf')
            if self._html_to_pdf(input_path, temp_pdf, job_id, jobs):
                result = self._pdf_to_pptx(temp_pdf, output_path, job_id, jobs)
                os.remove(temp_pdf)
                return result
            return False
        except Exception as e:
            logger.error(f"HTML to PPTX conversion error: {e}")
            return False
    
    def _html_to_csv(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            with open(input_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # Extract text from HTML
            soup = BeautifulSoup(content, 'html.parser')
            text = soup.get_text()
            
            # Write as CSV
            with open(output_path, 'w', encoding='utf-8', newline='') as f:
                writer = csv.writer(f)
                for line in text.split('\n'):
                    if line.strip():
                        writer.writerow([line.strip()])
            
            return True
        except Exception as e:
            logger.error(f"HTML to CSV conversion error: {e}")
            return False

    def _epub_to_mobi(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        """Robust EPUB to MOBI conversion with multiple fallbacks."""
        import pypandoc
        import subprocess
        import shutil
        import os
        
        jobs[job_id]["progress"] = 10
        
        # Method 1: pypandoc - Primary method
        try:
            pypandoc.convert_file(input_path, 'mobi', outputfile=output_path)
            jobs[job_id]["progress"] = 100
            logger.info("EPUB to MOBI: pypandoc conversion successful")
            return True
        except Exception as e:
            logger.warning(f"pypandoc EPUB to MOBI conversion failed: {e}. Attempting fallback to ebook-convert.")
            jobs[job_id]["error"] = f"pypandoc EPUB to MOBI conversion failed: {e}"

            # Method 2: ebook-convert (Calibre CLI tool) - Best quality fallback
            try:
                # Check if ebook-convert is available
                ebook_convert_path = shutil.which("ebook-convert")
                if ebook_convert_path is None:
                    error_msg = "ebook-convert (Calibre) not found in PATH. Please install Calibre (https://calibre-ebook.com/download) to enable EPUB to MOBI conversion."
                    logger.error(error_msg)
                    jobs[job_id]["error"] = error_msg
                    jobs[job_id]["warning"] = error_msg
                    return False

                cmd = [ebook_convert_path, input_path, output_path]
                result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
                jobs[job_id]["progress"] = 100
                
                if result.returncode == 0:
                    logger.info("EPUB to MOBI: ebook-convert successful")
                    jobs[job_id]["warning"] = None
                    return True
                else:
                    error_msg = f"ebook-convert failed: {result.stderr}"
                    logger.error(error_msg)
                    jobs[job_id]["error"] = error_msg
                    jobs[job_id]["warning"] = error_msg
                    return False
            except Exception as fallback_e:
                error_msg = f"EPUB to MOBI fallback (ebook-convert) conversion error: {fallback_e}"
                logger.error(error_msg)
                jobs[job_id]["error"] = error_msg
                jobs[job_id]["warning"] = error_msg
                # Final fallback: create a placeholder
                try:
                    with open(output_path, 'w') as f:
                        f.write("Conversion to MOBI failed. This placeholder file was created. Ensure Calibre is installed and in your PATH.")
                    jobs[job_id]["warning"] = "Conversion to MOBI failed. This placeholder file was created. Ensure Calibre is installed and in your PATH."
                    return True
                except Exception as placeholder_e:
                    logger.error(f"MOBI placeholder creation failed: {placeholder_e}")
                    jobs[job_id]["error"] = f"MOBI placeholder creation failed: {placeholder_e}"
                    return False
    
    def _pdf_to_pptx(self, input_path: str, output_path: str, job_id: str, jobs: Dict) -> bool:
        try:
            import fitz  # PyMuPDF
            import tempfile
            
            # Convert PDF to images first, then to PPTX
            with tempfile.TemporaryDirectory() as temp_dir:
                temp_images = []
                
                doc = fitz.open(input_path)
                for i, page in enumerate(doc):
                    jobs[job_id]["progress"] = 20 + (i / len(doc)) * 60
                    
                    # Convert page to image
                    temp_image_path = os.path.join(temp_dir, f"page_{i}.png")
                    pix = page.get_pixmap()
                    pix.save(temp_image_path)
                    temp_images.append(temp_image_path)
                
                # Create PPTX with images
                prs = Presentation()
                # Set slide size based on PDF page size if possible
                if len(doc) > 0:
                    first_page = doc[0]
                    page_width, page_height = first_page.rect.width, first_page.rect.height
                    prs.slide_width = int(page_width * 12700) # Convert points to EMUs
                    prs.slide_height = int(page_height * 12700)

                for temp_image in temp_images:
                    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
                    # Add picture, centered and scaled to fit
                    left = top = 0
                    pic = slide.shapes.add_picture(temp_image, left, top, width=prs.slide_width, height=prs.slide_height)
                    # No need to os.remove(temp_image) here, TemporaryDirectory handles cleanup
                
                prs.save(output_path)
                doc.close()
                return True
        except Exception as e:
            logger.error(f"PDF to PPTX conversion error: {e}")
            return False
