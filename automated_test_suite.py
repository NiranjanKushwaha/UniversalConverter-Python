#!/usr/bin/env python3
"""
Automated Test Suite for Universal File Converter
Tests all conversion combinations and generates detailed reports
"""

import os
import sys
import asyncio
import json
import time
import hashlib
import shutil
import tempfile
from pathlib import Path
from typing import Dict, List, Tuple, Optional
from datetime import datetime
import logging
import requests
from concurrent.futures import ThreadPoolExecutor, as_completed
import pandas as pd
from PIL import Image
import fitz  # PyMuPDF
from docx import Document
import openpyxl
from pptx import Presentation
import xml.etree.ElementTree as ET

# Setup logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

class ConversionTestSuite:
    def __init__(self, api_base_url: str = "http://localhost:8000", test_folder: str = "test_files"):
        self.api_base_url = api_base_url
        self.test_folder = Path(test_folder)
        self.results = []
        self.test_files = []
        self.supported_formats = {}
        
        # Create test directories
        self.test_output_dir = Path("test_outputs")
        self.test_output_dir.mkdir(exist_ok=True)
        
        # Content verification methods
        self.content_verifiers = {
            "PDF": self._verify_pdf_content,
            "DOCX": self._verify_docx_content,
            "XLSX": self._verify_xlsx_content,
            "PPTX": self._verify_pptx_content,
            "JPG": self._verify_image_content,
            "PNG": self._verify_image_content,
            "HTML": self._verify_html_content,
            "TXT": self._verify_txt_content,
        }
    
    def discover_test_files(self) -> List[Dict]:
        """Discover all test files in the test folder"""
        logger.info(f"Discovering test files in {self.test_folder}")
        
        if not self.test_folder.exists():
            logger.error(f"Test folder {self.test_folder} does not exist!")
            return []
        
        test_files = []
        supported_extensions = {
            '.pdf': 'PDF', '.docx': 'DOCX', '.doc': 'DOC', '.xlsx': 'XLSX', 
            '.xls': 'XLS', '.pptx': 'PPTX', '.ppt': 'PPT', '.jpg': 'JPG', 
            '.jpeg': 'JPG', '.png': 'PNG', '.gif': 'GIF', '.bmp': 'BMP', 
            '.html': 'HTML', '.htm': 'HTML', '.txt': 'TXT', '.csv': 'CSV',
            '.json': 'JSON', '.xml': 'XML', '.rtf': 'RTF', '.odt': 'ODT'
        }
        
        for file_path in self.test_folder.rglob("*"):
            if file_path.is_file():
                ext = file_path.suffix.lower()
                if ext in supported_extensions:
                    test_files.append({
                        'path': file_path,
                        'filename': file_path.name,
                        'format': supported_extensions[ext],
                        'size': file_path.stat().st_size,
                        'hash': self._calculate_file_hash(file_path)
                    })
        
        logger.info(f"Found {len(test_files)} test files")
        return test_files
    
    def _calculate_file_hash(self, file_path: Path) -> str:
        """Calculate SHA-256 hash of file"""
        hash_sha256 = hashlib.sha256()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_sha256.update(chunk)
        return hash_sha256.hexdigest()
    
    async def get_supported_formats(self) -> Dict:
        """Get supported conversion formats from API"""
        try:
            response = requests.get(f"{self.api_base_url}/formats")
            if response.status_code == 200:
                self.supported_formats = response.json()
                logger.info(f"Retrieved {len(self.supported_formats)} supported format combinations")
                return self.supported_formats
            else:
                logger.error(f"Failed to get supported formats: {response.status_code}")
                return {}
        except Exception as e:
            logger.error(f"Error getting supported formats: {e}")
            return {}
    
    def generate_test_combinations(self) -> List[Tuple]:
        """Generate all possible conversion combinations"""
        combinations = []
        
        for test_file in self.test_files:
            source_format = test_file['format']
            
            # Get possible destinations for this source format
            for format_info in self.supported_formats:
                if format_info['source'] == source_format:
                    for dest_format in format_info['destination']:
                        combinations.append((
                            test_file,
                            source_format,
                            dest_format
                        ))
        
        logger.info(f"Generated {len(combinations)} test combinations")
        return combinations
    
    async def test_conversion(self, test_file: Dict, source_format: str, dest_format: str) -> Dict:
        """Test a single conversion"""
        start_time = time.time()
        result = {
            'test_file': test_file['filename'],
            'source_format': source_format,
            'dest_format': dest_format,
            'start_time': datetime.now().isoformat(),
            'success': False,
            'error': None,
            'conversion_method': None,
            'warning': None,
            'duration': 0,
            'output_size': 0,
            'content_preserved': False,
            'content_verification': {}
        }
        
        try:
            # Upload and convert
            with open(test_file['path'], 'rb') as f:
                files = {'file': (test_file['filename'], f, 'application/octet-stream')}
                data = {
                    'sourceFormat': source_format,
                    'destinationFormat': dest_format
                }
                
                response = requests.post(f"{self.api_base_url}/convert", files=files, data=data)
                
                if response.status_code != 200:
                    result['error'] = f"Upload failed: {response.status_code}"
                    return result
                
                job_data = response.json()
                job_id = job_data['jobId']
                
                # Poll for completion
                max_wait = 300  # 5 minutes
                wait_time = 0
                while wait_time < max_wait:
                    status_response = requests.get(f"{self.api_base_url}/status/{job_id}")
                    if status_response.status_code == 200:
                        status_data = status_response.json()
                        
                        if status_data['status'] == 'completed':
                            result['success'] = True
                            result['conversion_method'] = status_data.get('conversion_method')
                            result['warning'] = status_data.get('warning')
                            
                            # Download the converted file
                            download_response = requests.get(f"{self.api_base_url}/download/{job_id}")
                            if download_response.status_code == 200:
                                output_filename = f"{test_file['filename']}_{source_format}_to_{dest_format}.{dest_format.lower()}"
                                output_path = self.test_output_dir / output_filename
                                
                                with open(output_path, 'wb') as f:
                                    f.write(download_response.content)
                                
                                result['output_size'] = output_path.stat().st_size
                                
                                # Verify content preservation
                                result['content_preserved'], result['content_verification'] = \
                                    self._verify_content_preservation(test_file['path'], output_path, source_format, dest_format)
                            
                            break
                        elif status_data['status'] == 'error':
                            result['error'] = status_data.get('error', 'Unknown error')
                            break
                    
                    await asyncio.sleep(2)
                    wait_time += 2
                
                if wait_time >= max_wait:
                    result['error'] = "Conversion timeout"
        
        except Exception as e:
            result['error'] = str(e)
        
        result['duration'] = time.time() - start_time
        return result
    
    def _verify_content_preservation(self, input_path: Path, output_path: Path, source_format: str, dest_format: str) -> Tuple[bool, Dict]:
        """Verify that content is preserved in conversion"""
        verification = {
            'text_preserved': False,
            'tables_preserved': False,
            'images_preserved': False,
            'structure_preserved': False,
            'details': {}
        }
        
        try:
            # Get verification method for source format
            source_verifier = self.content_verifiers.get(source_format)
            dest_verifier = self.content_verifiers.get(dest_format)
            
            if source_verifier and dest_verifier:
                source_content = source_verifier(input_path)
                dest_content = dest_verifier(output_path)
                
                # Compare content
                verification['text_preserved'] = self._compare_text_content(source_content, dest_content)
                verification['tables_preserved'] = self._compare_table_content(source_content, dest_content)
                verification['images_preserved'] = self._compare_image_content(source_content, dest_content)
                verification['structure_preserved'] = self._compare_structure(source_content, dest_content)
                verification['details'] = {
                    'source_content': source_content,
                    'dest_content': dest_content
                }
            
            # Overall preservation score
            preservation_score = sum([
                verification['text_preserved'],
                verification['tables_preserved'],
                verification['images_preserved'],
                verification['structure_preserved']
            ]) / 4.0
            
            return preservation_score > 0.5, verification
            
        except Exception as e:
            verification['details']['error'] = str(e)
            return False, verification
    
    def _verify_pdf_content(self, file_path: Path) -> Dict:
        """Extract content from PDF"""
        try:
            doc = fitz.open(str(file_path))
            content = {
                'text': '',
                'tables': [],
                'images': [],
                'pages': len(doc)
            }
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                content['text'] += page.get_text()
                
                # Extract images
                image_list = page.get_images()
                content['images'].extend([f"Page {page_num + 1}: {img[0]}" for img in image_list])
            
            doc.close()
            return content
        except Exception as e:
            return {'error': str(e)}
    
    def _verify_docx_content(self, file_path: Path) -> Dict:
        """Extract content from DOCX"""
        try:
            doc = Document(file_path)
            content = {
                'text': '',
                'tables': [],
                'images': [],
                'paragraphs': len(doc.paragraphs)
            }
            
            # Extract text
            for para in doc.paragraphs:
                content['text'] += para.text + '\n'
            
            # Extract tables
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                content['tables'].append(table_data)
            
            # Extract images (basic check)
            for para in doc.paragraphs:
                for run in para.runs:
                    if 'graphic' in run._element.xml:
                        content['images'].append(f"Image in paragraph: {para.text[:50]}...")
            
            return content
        except Exception as e:
            return {'error': str(e)}
    
    def _verify_xlsx_content(self, file_path: Path) -> Dict:
        """Extract content from XLSX"""
        try:
            wb = openpyxl.load_workbook(file_path)
            content = {
                'sheets': [],
                'tables': [],
                'text': ''
            }
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_data = []
                
                for row in sheet.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        sheet_data.append([str(cell) if cell is not None else '' for cell in row])
                        content['text'] += ' '.join([str(cell) if cell is not None else '' for cell in row]) + '\n'
                
                if sheet_data:
                    content['sheets'].append(sheet_data)
                    content['tables'].append(sheet_data)
            
            return content
        except Exception as e:
            return {'error': str(e)}
    
    def _verify_pptx_content(self, file_path: Path) -> Dict:
        """Extract content from PPTX"""
        try:
            prs = Presentation(file_path)
            content = {
                'slides': [],
                'text': '',
                'images': []
            }
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = f"Slide {slide_num + 1}: "
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + " "
                        content['text'] += shape.text + "\n"
                
                content['slides'].append(slide_text)
            
            return content
        except Exception as e:
            return {'error': str(e)}
    
    def _verify_image_content(self, file_path: Path) -> Dict:
        """Extract content from image"""
        try:
            with Image.open(file_path) as img:
                return {
                    'size': img.size,
                    'mode': img.mode,
                    'format': img.format,
                    'text': f"Image: {img.size[0]}x{img.size[1]} {img.mode}"
                }
        except Exception as e:
            return {'error': str(e)}
    
    def _verify_html_content(self, file_path: Path) -> Dict:
        """Extract content from HTML"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            return {
                'text': content,
                'length': len(content),
                'has_tables': '<table' in content.lower(),
                'has_images': '<img' in content.lower()
            }
        except Exception as e:
            return {'error': str(e)}
    
    def _verify_txt_content(self, file_path: Path) -> Dict:
        """Extract content from TXT"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            return {
                'text': content,
                'length': len(content),
                'lines': len(content.split('\n'))
            }
        except Exception as e:
            return {'error': str(e)}
    
    def _compare_text_content(self, source_content: Dict, dest_content: Dict) -> bool:
        """Compare text content between source and destination"""
        try:
            source_text = source_content.get('text', '').lower()
            dest_text = dest_content.get('text', '').lower()
            
            # Remove whitespace and special characters for comparison
            source_clean = ''.join(c for c in source_text if c.isalnum() or c.isspace())
            dest_clean = ''.join(c for c in dest_text if c.isalnum() or c.isspace())
            
            # Calculate similarity
            if len(source_clean) == 0:
                return len(dest_clean) == 0
            
            similarity = len(set(source_clean.split()) & set(dest_clean.split())) / len(set(source_clean.split()))
            return similarity > 0.7  # 70% similarity threshold
        except:
            return False
    
    def _compare_table_content(self, source_content: Dict, dest_content: Dict) -> bool:
        """Compare table content between source and destination"""
        try:
            source_tables = source_content.get('tables', [])
            dest_tables = dest_content.get('tables', [])
            
            if len(source_tables) == 0:
                return len(dest_tables) == 0
            
            # Check if tables are preserved
            return len(dest_tables) >= len(source_tables) * 0.8  # 80% of tables preserved
        except:
            return False
    
    def _compare_image_content(self, source_content: Dict, dest_content: Dict) -> bool:
        """Compare image content between source and destination"""
        try:
            source_images = source_content.get('images', [])
            dest_images = dest_content.get('images', [])
            
            if len(source_images) == 0:
                return len(dest_images) == 0
            
            # Check if images are preserved
            return len(dest_images) >= len(source_images) * 0.8  # 80% of images preserved
        except:
            return False
    
    def _compare_structure(self, source_content: Dict, dest_content: Dict) -> bool:
        """Compare overall structure between source and destination"""
        try:
            # Compare basic structure elements
            source_keys = set(source_content.keys())
            dest_keys = set(dest_content.keys())
            
            # Check if key structural elements are preserved
            structural_elements = {'text', 'tables', 'images', 'pages', 'slides', 'sheets'}
            preserved_elements = source_keys & dest_keys & structural_elements
            
            return len(preserved_elements) >= len(source_keys & structural_elements) * 0.7
        except:
            return False
    
    async def run_full_test_suite(self) -> Dict:
        """Run the complete test suite"""
        logger.info("Starting comprehensive test suite...")
        
        # Discover test files
        self.test_files = self.discover_test_files()
        if not self.test_files:
            logger.error("No test files found!")
            return {}
        
        # Get supported formats
        await self.get_supported_formats()
        if not self.supported_formats:
            logger.error("Could not retrieve supported formats!")
            return {}
        
        # Generate test combinations
        combinations = self.generate_test_combinations()
        if not combinations:
            logger.error("No test combinations generated!")
            return {}
        
        # Run tests
        logger.info(f"Running {len(combinations)} conversion tests...")
        
        # Use ThreadPoolExecutor for concurrent testing
        with ThreadPoolExecutor(max_workers=5) as executor:
            futures = []
            for test_file, source_format, dest_format in combinations:
                future = executor.submit(
                    asyncio.run, 
                    self.test_conversion(test_file, source_format, dest_format)
                )
                futures.append(future)
            
            # Collect results
            for future in as_completed(futures):
                try:
                    result = future.result()
                    self.results.append(result)
                    logger.info(f"Completed: {result['test_file']} {result['source_format']} → {result['dest_format']} - {'SUCCESS' if result['success'] else 'FAILED'}")
                except Exception as e:
                    logger.error(f"Test failed with exception: {e}")
        
        # Generate report
        return self.generate_report()
    
    def generate_report(self) -> Dict:
        """Generate comprehensive test report"""
        logger.info("Generating test report...")
        
        if not self.results:
            return {}
        
        # Calculate statistics
        total_tests = len(self.results)
        successful_tests = sum(1 for r in self.results if r['success'])
        failed_tests = total_tests - successful_tests
        
        # Group by format combinations
        format_stats = {}
        for result in self.results:
            key = f"{result['source_format']} → {result['dest_format']}"
            if key not in format_stats:
                format_stats[key] = {'total': 0, 'success': 0, 'fail': 0, 'avg_duration': 0}
            
            format_stats[key]['total'] += 1
            if result['success']:
                format_stats[key]['success'] += 1
            else:
                format_stats[key]['fail'] += 1
        
        # Calculate average duration for each format
        for key in format_stats:
            durations = [r['duration'] for r in self.results 
                        if f"{r['source_format']} → {r['dest_format']}" == key and r['success']]
            if durations:
                format_stats[key]['avg_duration'] = sum(durations) / len(durations)
        
        # Content preservation analysis
        content_preservation_stats = {
            'text_preserved': sum(1 for r in self.results if r.get('content_verification', {}).get('text_preserved', False)),
            'tables_preserved': sum(1 for r in self.results if r.get('content_verification', {}).get('tables_preserved', False)),
            'images_preserved': sum(1 for r in self.results if r.get('content_verification', {}).get('images_preserved', False)),
            'structure_preserved': sum(1 for r in self.results if r.get('content_verification', {}).get('structure_preserved', False))
        }
        
        # Error analysis
        error_analysis = {}
        for result in self.results:
            if result['error']:
                error_type = result['error'].split(':')[0] if ':' in result['error'] else result['error']
                error_analysis[error_type] = error_analysis.get(error_type, 0) + 1
        
        # Generate report
        report = {
            'summary': {
                'total_tests': total_tests,
                'successful_tests': successful_tests,
                'failed_tests': failed_tests,
                'success_rate': (successful_tests / total_tests) * 100 if total_tests > 0 else 0,
                'test_duration': sum(r['duration'] for r in self.results),
                'timestamp': datetime.now().isoformat()
            },
            'format_statistics': format_stats,
            'content_preservation': content_preservation_stats,
            'error_analysis': error_analysis,
            'detailed_results': self.results,
            'recommendations': self._generate_recommendations()
        }
        
        # Save report
        report_path = self.test_output_dir / f"test_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json"
        with open(report_path, 'w') as f:
            json.dump(report, f, indent=2)
        
        # Generate CSV summary
        self._generate_csv_summary(report)
        
        # Generate HTML report
        self._generate_html_report(report)
        
        logger.info(f"Test report saved to {report_path}")
        return report
    
    def _generate_recommendations(self) -> List[str]:
        """Generate recommendations based on test results"""
        recommendations = []
        
        # Analyze success rates
        success_rate = sum(1 for r in self.results if r['success']) / len(self.results) * 100
        
        if success_rate < 80:
            recommendations.append("Overall success rate is low. Consider improving error handling and fallback mechanisms.")
        
        # Analyze content preservation
        content_preserved = sum(1 for r in self.results if r.get('content_preserved', False))
        content_rate = content_preserved / len(self.results) * 100
        
        if content_rate < 70:
            recommendations.append("Content preservation rate is low. Consider improving conversion methods for better fidelity.")
        
        # Analyze specific format issues
        format_issues = {}
        for result in self.results:
            if not result['success']:
                key = f"{result['source_format']} → {result['dest_format']}"
                format_issues[key] = format_issues.get(key, 0) + 1
        
        for format_combo, count in format_issues.items():
            if count > 2:
                recommendations.append(f"Multiple failures for {format_combo}. Consider implementing better conversion methods.")
        
        # Check for missing conversion methods
        missing_methods = []
        for result in self.results:
            if result['success'] and result.get('conversion_method') == 'python-docx-fallback':
                missing_methods.append(f"{result['source_format']} → {result['dest_format']}")
        
        if missing_methods:
            recommendations.append(f"Fallback methods used for: {', '.join(set(missing_methods))}. Consider installing LibreOffice or other conversion tools.")
        
        return recommendations
    
    def _generate_csv_summary(self, report: Dict):
        """Generate CSV summary of test results"""
        csv_data = []
        
        for result in report['detailed_results']:
            csv_data.append({
                'Test File': result['test_file'],
                'Source Format': result['source_format'],
                'Destination Format': result['dest_format'],
                'Success': result['success'],
                'Duration (s)': round(result['duration'], 2),
                'Conversion Method': result.get('conversion_method', 'N/A'),
                'Warning': result.get('warning', 'N/A'),
                'Error': result.get('error', 'N/A'),
                'Content Preserved': result.get('content_preserved', False),
                'Output Size (bytes)': result.get('output_size', 0)
            })
        
        df = pd.DataFrame(csv_data)
        csv_path = self.test_output_dir / f"test_summary_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv"
        df.to_csv(csv_path, index=False)
        logger.info(f"CSV summary saved to {csv_path}")
    
    def _generate_html_report(self, report: Dict):
        """Generate HTML report for better visualization"""
        html_content = f"""
        <!DOCTYPE html>
        <html>
        <head>
            <title>Universal File Converter Test Report</title>
            <style>
                body {{ font-family: Arial, sans-serif; margin: 20px; }}
                .summary {{ background: #f5f5f5; padding: 20px; border-radius: 5px; }}
                .stats {{ display: grid; grid-template-columns: repeat(auto-fit, minmax(200px, 1fr)); gap: 20px; margin: 20px 0; }}
                .stat-card {{ background: white; padding: 15px; border-radius: 5px; box-shadow: 0 2px 4px rgba(0,0,0,0.1); }}
                .success {{ color: green; }}
                .failure {{ color: red; }}
                table {{ width: 100%; border-collapse: collapse; margin: 20px 0; }}
                th, td {{ border: 1px solid #ddd; padding: 8px; text-align: left; }}
                th {{ background-color: #f2f2f2; }}
                .recommendations {{ background: #fff3cd; padding: 15px; border-radius: 5px; margin: 20px 0; }}
            </style>
        </head>
        <body>
            <h1>Universal File Converter Test Report</h1>
            <div class="summary">
                <h2>Test Summary</h2>
                <div class="stats">
                    <div class="stat-card">
                        <h3>Total Tests</h3>
                        <p>{report['summary']['total_tests']}</p>
                    </div>
                    <div class="stat-card">
                        <h3>Successful</h3>
                        <p class="success">{report['summary']['successful_tests']}</p>
                    </div>
                    <div class="stat-card">
                        <h3>Failed</h3>
                        <p class="failure">{report['summary']['failed_tests']}</p>
                    </div>
                    <div class="stat-card">
                        <h3>Success Rate</h3>
                        <p>{report['summary']['success_rate']:.1f}%</p>
                    </div>
                </div>
            </div>
            
            <h2>Format Statistics</h2>
            <table>
                <tr>
                    <th>Format Combination</th>
                    <th>Total</th>
                    <th>Success</th>
                    <th>Fail</th>
                    <th>Success Rate</th>
                    <th>Avg Duration (s)</th>
                </tr>
        """
        
        for format_combo, stats in report['format_statistics'].items():
            success_rate = (stats['success'] / stats['total']) * 100 if stats['total'] > 0 else 0
            html_content += f"""
                <tr>
                    <td>{format_combo}</td>
                    <td>{stats['total']}</td>
                    <td class="success">{stats['success']}</td>
                    <td class="failure">{stats['fail']}</td>
                    <td>{success_rate:.1f}%</td>
                    <td>{stats['avg_duration']:.2f}</td>
                </tr>
            """
        
        html_content += """
            </table>
            
            <h2>Content Preservation Analysis</h2>
            <div class="stats">
        """
        
        for content_type, count in report['content_preservation'].items():
            percentage = (count / report['summary']['total_tests']) * 100 if report['summary']['total_tests'] > 0 else 0
            html_content += f"""
                <div class="stat-card">
                    <h3>{content_type.replace('_', ' ').title()}</h3>
                    <p>{count} ({percentage:.1f}%)</p>
                </div>
            """
        
        html_content += """
            </div>
            
            <h2>Error Analysis</h2>
            <table>
                <tr>
                    <th>Error Type</th>
                    <th>Count</th>
                </tr>
        """
        
        for error_type, count in report['error_analysis'].items():
            html_content += f"""
                <tr>
                    <td>{error_type}</td>
                    <td>{count}</td>
                </tr>
            """
        
        html_content += """
            </table>
            
            <h2>Recommendations</h2>
            <div class="recommendations">
        """
        
        for recommendation in report['recommendations']:
            html_content += f"<p>• {recommendation}</p>"
        
        html_content += """
            </div>
            
            <h2>Detailed Results</h2>
            <table>
                <tr>
                    <th>Test File</th>
                    <th>Source → Dest</th>
                    <th>Status</th>
                    <th>Method</th>
                    <th>Duration</th>
                    <th>Content Preserved</th>
                </tr>
        """
        
        for result in report['detailed_results']:
            status_class = "success" if result['success'] else "failure"
            status_text = "SUCCESS" if result['success'] else "FAILED"
            html_content += f"""
                <tr>
                    <td>{result['test_file']}</td>
                    <td>{result['source_format']} → {result['dest_format']}</td>
                    <td class="{status_class}">{status_text}</td>
                    <td>{result.get('conversion_method', 'N/A')}</td>
                    <td>{result['duration']:.2f}s</td>
                    <td>{'Yes' if result.get('content_preserved', False) else 'No'}</td>
                </tr>
            """
        
        html_content += """
            </table>
        </body>
        </html>
        """
        
        html_path = self.test_output_dir / f"test_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.html"
        with open(html_path, 'w') as f:
            f.write(html_content)
        
        logger.info(f"HTML report saved to {html_path}")

async def main():
    """Main function to run the test suite"""
    import argparse
    
    parser = argparse.ArgumentParser(description='Universal File Converter Test Suite')
    parser.add_argument('--api-url', default='http://localhost:8000', help='API base URL')
    parser.add_argument('--test-folder', default='test_files', help='Folder containing test files')
    parser.add_argument('--output-dir', default='test_outputs', help='Output directory for results')
    
    args = parser.parse_args()
    
    # Create test suite
    test_suite = ConversionTestSuite(args.api_url, args.test_folder)
    
    # Run tests
    report = await test_suite.run_full_test_suite()
    
    if report:
        print(f"\n{'='*60}")
        print("TEST SUITE COMPLETED")
        print(f"{'='*60}")
        print(f"Total Tests: {report['summary']['total_tests']}")
        print(f"Successful: {report['summary']['successful_tests']}")
        print(f"Failed: {report['summary']['failed_tests']}")
        print(f"Success Rate: {report['summary']['success_rate']:.1f}%")
        print(f"Total Duration: {report['summary']['test_duration']:.2f} seconds")
        print(f"\nReports saved to: {test_suite.test_output_dir}")
        print(f"{'='*60}")
        
        if report['recommendations']:
            print("\nRECOMMENDATIONS:")
            for rec in report['recommendations']:
                print(f"• {rec}")
    else:
        print("Test suite failed to run!")

if __name__ == "__main__":
    asyncio.run(main()) 